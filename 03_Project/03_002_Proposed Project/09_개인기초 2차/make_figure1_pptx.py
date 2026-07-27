"""
그림 1 - 연구 개념도 PPTX 생성 스크립트
(좌) 문제: Cu 위 불균일 Li 핵생성 및 덴드라이트
(중) ALD ZnF2 공정 사이클 및 박막 형성
(우) 해결: 불화물 인공 SEI → 균일 Li 증착
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── 색상 팔레트 ──────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

C_HEADER      = rgb(0x28, 0x35, 0x93)
C_WHITE       = rgb(0xFF, 0xFF, 0xFF)
C_BLACK       = rgb(0x21, 0x21, 0x21)
C_CU          = rgb(0xB8, 0x73, 0x33)   # 구리
C_LI          = rgb(0x62, 0x62, 0x62)   # Li 금속 (어두운 회색)
C_LI_FLAT     = rgb(0x90, 0x90, 0x90)   # 균일 Li 층 (중간 회색)
C_ELECTROLYTE = rgb(0xE1, 0xF5, 0xFE)   # 전해질 (연파랑)
C_SEI_BAD     = rgb(0xFF, 0x5F, 0x22)   # 불균일 SEI (주황)
C_ALD_FILM    = rgb(0x00, 0x96, 0x88)   # ALD 불화물 박막 (틸)
C_PROBLEM     = rgb(0xC6, 0x28, 0x28)   # 문제 (적색)
C_SOLUTION    = rgb(0x2E, 0x7D, 0x32)   # 해결 (녹색)
C_PROCESS     = rgb(0x30, 0x3F, 0x9F)   # ALD 공정 (남색)
C_STEP1       = rgb(0x19, 0x76, 0xD2)   # 전구체 (파랑)
C_STEP2       = rgb(0xF9, 0xA8, 0x25)   # 퍼지1 (황색)
C_STEP3       = rgb(0x43, 0xA0, 0x47)   # 반응물 (초록)
C_STEP4       = rgb(0xEF, 0x53, 0x50)   # 퍼지2 (빨강)
C_ARROW       = rgb(0x37, 0x47, 0x4F)
C_GRID        = rgb(0xBD, 0xBD, 0xBD)
C_PLACEHOLDER = rgb(0xEE, 0xEE, 0xEE)
C_PH_TEXT     = rgb(0x75, 0x75, 0x75)

# ── 슬라이드 크기 ─────────────────────────────────────────────
W = 25.4   # cm  (A4 가로)
H = 15.5   # cm

prs = Presentation()
prs.slide_width  = Cm(W)
prs.slide_height = Cm(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── 도형 헬퍼 ─────────────────────────────────────────────────
def rect(x, y, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def rrect(x, y, w, h, fill=None, line=None, lw=Pt(0)):
    s = slide.shapes.add_shape(5, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def triangle(x, y, w, h, fill=C_LI):
    s = slide.shapes.add_shape(7, Cm(x), Cm(y), Cm(w), Cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def circle(x, y, d, fill=C_LI):
    s = slide.shapes.add_shape(9, Cm(x), Cm(y), Cm(d), Cm(d))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s

def tb(text, x, y, w, h, font="맑은 고딕", size=9, bold=False,
       color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def mtb(lines, x, y, w, font="맑은 고딕", size=8.5,
        color=C_BLACK, lsp=13.0):
    h_est = len(lines) * lsp / 28.35 + 0.4
    txb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h_est))
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = Pt(lsp)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color

def line(x1, y1, x2, y2, color=C_ARROW, width=Pt(1.5)):
    conn = slide.shapes.add_connector(1, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    conn.line.color.rgb = color
    conn.line.width = width
    return conn

def arrow(x1, y1, x2, y2, color=C_ARROW, width=Pt(2.0)):
    conn = slide.shapes.add_connector(1, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    conn.line.color.rgb = color
    conn.line.width = width
    # tailEnd = arrowhead at (x2, y2)
    tail_xml = (
        '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' type="arrow" w="med" len="med"/>'
    )
    conn.line._ln.append(etree.fromstring(tail_xml))
    return conn


# ═══════════════════════════════════════════════════════════════
# § 1. 헤더
# ═══════════════════════════════════════════════════════════════
rect(0, 0, W, 1.0, fill=C_HEADER)
tb("[그림 1]  연구 개념도 — ALD 불화물 인공 SEI 전략",
   0, 0.06, W, 0.88, size=13, bold=True, color=C_WHITE,
   align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# § 2. 패널 레이아웃
# ═══════════════════════════════════════════════════════════════
PANEL_Y = 1.05
PANEL_H = 14.3

# 패널 폭: (25.4 - 0.25 - 0.25 - 0.7 - 0.7) / 3 = 23.5 / 3 ≈ 7.83
# 패널 1: x=0.25,     w=7.7
# 갭 1:   x=7.95,     w=0.7 (화살표)
# 패널 2: x=8.65,     w=7.7
# 갭 2:   x=16.35,    w=0.7
# 패널 3: x=17.05,    w=8.1
P1_X  = 0.25;  PW  = 7.7
GAP1  = 7.95;  GW  = 0.7
P2_X  = 8.65
GAP2  = 16.35
P3_X  = 17.05; P3W = 8.1

# 패널 배경 및 테두리
rect(P1_X, PANEL_Y, PW,  PANEL_H, fill=C_WHITE, line=C_PROBLEM,  lw=Pt(1.5))
rect(P2_X, PANEL_Y, PW,  PANEL_H, fill=C_WHITE, line=C_PROCESS,  lw=Pt(1.5))
rect(P3_X, PANEL_Y, P3W, PANEL_H, fill=C_WHITE, line=C_SOLUTION, lw=Pt(1.5))

# 타이틀 바
TITLE_H = 0.75
rect(P1_X, PANEL_Y, PW,  TITLE_H, fill=C_PROBLEM)
rect(P2_X, PANEL_Y, PW,  TITLE_H, fill=C_PROCESS)
rect(P3_X, PANEL_Y, P3W, TITLE_H, fill=C_SOLUTION)

tb("(좌)  문제 : 불균일 Li 핵생성 및 덴드라이트",
   P1_X, PANEL_Y + 0.06, PW, 0.63, size=9, bold=True,
   color=C_WHITE, align=PP_ALIGN.CENTER)
tb("(중)  ALD ZnF₂ 증착 공정 사이클",
   P2_X, PANEL_Y + 0.06, PW, 0.63, size=9, bold=True,
   color=C_WHITE, align=PP_ALIGN.CENTER)
tb("(우)  해결 : 균일 Li 핵생성 및 평탄한 증착",
   P3_X, PANEL_Y + 0.06, P3W, 0.63, size=9, bold=True,
   color=C_WHITE, align=PP_ALIGN.CENTER)

# ── 패널 간 화살표 ──────────────────────────────────────────
ARR_Y  = PANEL_Y + PANEL_H * 0.38
ARR_Y2 = PANEL_Y + PANEL_H * 0.38

arrow(GAP1 + 0.05, ARR_Y, GAP1 + GW - 0.05, ARR_Y,
      color=C_PROCESS, width=Pt(3))
tb("ALD\n코팅", GAP1, ARR_Y - 0.60, GW, 0.55,
   size=7.5, bold=True, color=C_PROCESS, align=PP_ALIGN.CENTER)

arrow(GAP2 + 0.05, ARR_Y2, GAP2 + GW - 0.05, ARR_Y2,
      color=C_SOLUTION, width=Pt(3))
tb("→ 해결", GAP2, ARR_Y2 - 0.60, GW, 0.55,
   size=7.5, bold=True, color=C_SOLUTION, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# § 3. 패널 1 — 문제: 불균일 Li 핵생성 및 덴드라이트
# ═══════════════════════════════════════════════════════════════
SCH_Y   = PANEL_Y + TITLE_H + 0.05   # 모식도 시작 y
SCH_H   = 6.8                         # 모식도 높이
SCH_W1  = PW - 0.10
CU_TOP  = SCH_Y + SCH_H - 0.05       # Cu 집전체 상단 y

# 전해질 배경
rect(P1_X + 0.05, SCH_Y, SCH_W1, SCH_H - 0.85, fill=C_ELECTROLYTE)
tb("전해질 (Electrolyte)", P1_X + 0.15, SCH_Y + 0.08,
   SCH_W1, 0.42, size=8, bold=True,
   color=rgb(0x01, 0x57, 0x9B))

# Cu 집전체
rect(P1_X + 0.05, CU_TOP - 0.85, SCH_W1, 0.85, fill=C_CU)
tb("Cu 집전체", P1_X + 0.05, CU_TOP - 0.65,
   SCH_W1, 0.45, size=9, bold=True,
   color=C_WHITE, align=PP_ALIGN.CENTER)

# 덴드라이트 (삼각형들 - 다양한 높이)
DEND = [
    # (x_center, height, width)  ← Cu_TOP 기준
    (P1_X + 0.35, 3.6, 0.33), (P1_X + 0.80, 2.3, 0.30),
    (P1_X + 1.25, 4.4, 0.42), (P1_X + 1.75, 1.8, 0.26),
    (P1_X + 2.20, 3.0, 0.31), (P1_X + 2.70, 4.8, 0.45),
    (P1_X + 3.20, 2.1, 0.28), (P1_X + 3.70, 3.7, 0.36),
    (P1_X + 4.20, 1.6, 0.24), (P1_X + 4.70, 4.1, 0.40),
    (P1_X + 5.20, 2.7, 0.32), (P1_X + 5.70, 3.3, 0.34),
    (P1_X + 6.20, 4.5, 0.43), (P1_X + 6.70, 2.0, 0.27),
]
for cx, dh, dw in DEND:
    dy = CU_TOP - 0.85 - dh
    triangle(cx - dw / 2, dy, dw, dh, fill=C_LI)

# 불균일 SEI 조각들 (orange segments, 불연속적으로)
SEI_Y = CU_TOP - 0.85 - 0.18
for sx, sw in [
    (P1_X+0.10, 0.50), (P1_X+0.75, 0.55), (P1_X+1.50, 0.48),
    (P1_X+2.25, 0.52), (P1_X+3.05, 0.50), (P1_X+3.80, 0.55),
    (P1_X+4.55, 0.48), (P1_X+5.30, 0.52), (P1_X+6.05, 0.50),
    (P1_X+6.75, 0.42),
]:
    rrect(sx, SEI_Y, sw, 0.20, fill=C_SEI_BAD)

# 라벨 (화살표 + 텍스트)
tb("Li 덴드라이트 ⚠", P1_X + 3.3, SCH_Y + 0.55, 3.8, 0.50,
   size=9, bold=True, color=C_PROBLEM)
tb("불균일·파단 SEI", P1_X + 3.8, SEI_Y - 0.55, 3.2, 0.45,
   size=8.5, color=rgb(0xBF, 0x36, 0x0C))

# § 설명 텍스트
DESC1_Y = SCH_Y + SCH_H + 0.10
rect(P1_X + 0.05, DESC1_Y, SCH_W1, 2.35,
     fill=rgb(0xFF, 0xEB, 0xEE), line=C_PROBLEM, lw=Pt(0.5))
mtb([
    "• Cu 집전체 위 ALD 코팅 없이 Li 전착 시",
    "• 핵생성 과전압 불균일 → 덴드라이트 성장",
    "• 두껍고 불균일한 SEI 형성 반복",
    "• 쿨롱 효율 저하 및 내부 단락(short) 위험",
], P1_X + 0.15, DESC1_Y + 0.10, SCH_W1 - 0.25,
   color=rgb(0xB7, 0x1C, 0x1C), size=8.5, lsp=13.0)

# § SEM 이미지 플레이스홀더
PH1_Y  = DESC1_Y + 2.50
PH1_H  = PANEL_H - (PH1_Y - PANEL_Y) - 0.12
rect(P1_X + 0.05, PH1_Y, SCH_W1, PH1_H,
     fill=C_PLACEHOLDER, line=C_GRID, lw=Pt(0.75))
# 대각선 (이미지 플레이스홀더 표시)
line(P1_X + 0.05, PH1_Y, P1_X + 0.05 + SCH_W1, PH1_Y + PH1_H,
     color=C_GRID, width=Pt(0.5))
line(P1_X + 0.05 + SCH_W1, PH1_Y, P1_X + 0.05, PH1_Y + PH1_H,
     color=C_GRID, width=Pt(0.5))
tb("SEM 이미지 공간\n(Li 덴드라이트 / 불균일 증착)",
   P1_X + 0.05, PH1_Y + PH1_H / 2 - 0.35, SCH_W1, 0.80,
   size=8, color=C_PH_TEXT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# § 4. 패널 2 — ALD ZnF2 공정 사이클
# ═══════════════════════════════════════════════════════════════
SCH_W2  = PW - 0.10
CY_Y    = PANEL_Y + TITLE_H + 0.12   # 사이클 다이어그램 시작
CY_H    = 5.6
CY_X    = P2_X + 0.05

# 사이클 다이어그램 배경
rect(CY_X, CY_Y, SCH_W2, CY_H, fill=rgb(0xF3, 0xE5, 0xF5),
     line=rgb(0x9C, 0x27, 0xB0), lw=Pt(0.5))
tb("ALD 증착 사이클 (1 cycle)", CY_X, CY_Y + 0.07,
   SCH_W2, 0.40, size=8.5, bold=True,
   color=rgb(0x6A, 0x1B, 0x9A), align=PP_ALIGN.CENTER)

# 4 스텝 박스 위치
SB_W = 2.85; SB_H = 1.05
GAP_X = (SCH_W2 - 2 * SB_W) / 3   # ≈ 0.60 cm
GAP_Y = (CY_H - 0.55 - 2 * SB_H) / 3  # ≈ 1.0 cm

s1x = CY_X + GAP_X
s1y = CY_Y + 0.55 + GAP_Y
s2x = CY_X + GAP_X * 2 + SB_W
s2y = s1y
s3x = s2x
s3y = CY_Y + 0.55 + GAP_Y * 2 + SB_H
s4x = s1x
s4y = s3y

STEPS = [
    (s1x, s1y, "① 전구체 노출", "(DEZ 펄스)", C_STEP1),
    (s2x, s2y, "② 퍼지", "(N₂ 가스)",   C_STEP2),
    (s3x, s3y, "③ 반응물 노출", "(HF·py 펄스)", C_STEP3),
    (s4x, s4y, "④ 퍼지", "(N₂ 가스)",   C_STEP4),
]
for sx, sy, st, sd, sc in STEPS:
    rrect(sx, sy, SB_W, SB_H, fill=sc, line=C_WHITE, lw=Pt(1.0))
    tb(st, sx, sy + 0.10, SB_W, 0.50,
       size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tb(sd, sx, sy + 0.58, SB_W, 0.42,
       size=8, color=C_WHITE, align=PP_ALIGN.CENTER)

# 스텝 간 화살표
def step_arrow(x1, y1, x2, y2):
    conn = slide.shapes.add_connector(1, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    conn.line.color.rgb = C_ARROW
    conn.line.width = Pt(2.0)
    tail_xml = (
        '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' type="arrow" w="med" len="med"/>'
    )
    conn.line._ln.append(etree.fromstring(tail_xml))

# 1→2 (오른쪽), 2→3 (아래), 3→4 (왼쪽), 4→1 (위)
step_arrow(s1x + SB_W,     s1y + SB_H / 2, s2x,           s2y + SB_H / 2)
step_arrow(s2x + SB_W / 2, s2y + SB_H,     s3x + SB_W / 2, s3y)
step_arrow(s3x,             s3y + SB_H / 2, s4x + SB_W,   s4y + SB_H / 2)
step_arrow(s4x + SB_W / 2, s4y,             s1x + SB_W / 2, s1y + SB_H)

# 중앙 라벨
CY_CX = CY_X + SCH_W2 / 2   # 사이클 다이어그램 중심 x
CY_CY = (s1y + SB_H + s3y) / 2  # 상하 박스 사이 중간
rrect(CY_CX - 1.10, CY_CY - 0.30, 2.20, 0.60,
      fill=C_WHITE, line=rgb(0x9C, 0x27, 0xB0), lw=Pt(1.0))
tb("ALD 1 사이클 반복", CY_CX - 1.10, CY_CY - 0.30,
   2.20, 0.60, size=8, bold=True,
   color=rgb(0x6A, 0x1B, 0x9A), align=PP_ALIGN.CENTER)

# § 단면도 (박막 성장)
XS_Y  = CY_Y + CY_H + 0.15
XS_H  = 1.85
XS_W  = SCH_W2

rect(CY_X, XS_Y, XS_W, XS_H, fill=rgb(0xE8, 0xEA, 0xF6),
     line=rgb(0x7B, 0x86, 0xCB), lw=Pt(0.5))
tb("박막 형성 단면 모식도", CY_X, XS_Y + 0.06,
   XS_W, 0.40, size=8, bold=True,
   color=C_PROCESS, align=PP_ALIGN.CENTER)

# Cu 기판
CU2_Y = XS_Y + 1.25
rect(CY_X + 0.25, CU2_Y, XS_W - 0.50, 0.42, fill=C_CU)
tb("Cu 기판 / Li 금속 음극", CY_X + 0.25, CU2_Y + 0.06,
   XS_W - 0.50, 0.32, size=8, bold=True,
   color=C_WHITE, align=PP_ALIGN.CENTER)

# ALD ZnF2 박막
ALD2_Y = CU2_Y - 0.28
rect(CY_X + 0.25, ALD2_Y, XS_W - 0.50, 0.28, fill=C_ALD_FILM)
tb("ZnF₂ ALD 박막  (두께: 5~100 nm)", CY_X + 0.25, ALD2_Y + 0.02,
   XS_W - 0.50, 0.28, size=7.5, bold=True,
   color=C_WHITE, align=PP_ALIGN.CENTER)

# 두께 화살표 주석
tb("← n 사이클 반복으로 두께 제어 →", CY_X + 0.25, ALD2_Y - 0.38,
   XS_W - 0.50, 0.33, size=7.5,
   color=rgb(0x39, 0x49, 0xAB), align=PP_ALIGN.CENTER)

# § 설명 텍스트
DESC2_Y = XS_Y + XS_H + 0.10
rect(CY_X, DESC2_Y, SCH_W2, 2.35,
     fill=rgb(0xEE, 0xEE, 0xFF), line=C_PROCESS, lw=Pt(0.5))
mtb([
    "• ZnF₂ ALD: DEZ + HF·py 전구체 조합",
    "• 자기제한적(self-limiting) 표면 반응",
    "• 기판 온도 50~200°C, Å 수준 두께 제어",
    "• 결과: 균일·공형(conformal) 불화물 박막",
], CY_X + 0.12, DESC2_Y + 0.10, SCH_W2 - 0.20,
   color=rgb(0x1A, 0x23, 0x7E), size=8.5, lsp=13.0)

# § TEM 이미지 플레이스홀더
PH2_Y  = DESC2_Y + 2.50
PH2_H  = PANEL_H - (PH2_Y - PANEL_Y) - 0.12
rect(CY_X, PH2_Y, SCH_W2, PH2_H,
     fill=C_PLACEHOLDER, line=C_GRID, lw=Pt(0.75))
line(CY_X, PH2_Y, CY_X + SCH_W2, PH2_Y + PH2_H,
     color=C_GRID, width=Pt(0.5))
line(CY_X + SCH_W2, PH2_Y, CY_X, PH2_Y + PH2_H,
     color=C_GRID, width=Pt(0.5))
tb("TEM / SAED 이미지 공간\n(ALD ZnF₂ 박막 단면)",
   CY_X, PH2_Y + PH2_H / 2 - 0.35, SCH_W2, 0.80,
   size=8, color=C_PH_TEXT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# § 5. 패널 3 — 해결: 균일 Li 핵생성 및 평탄한 증착
# ═══════════════════════════════════════════════════════════════
SCH_W3 = P3W - 0.10

# 전해질 배경 (Panel 1과 동일한 y 범위)
rect(P3_X + 0.05, SCH_Y, SCH_W3, SCH_H - 0.85, fill=C_ELECTROLYTE)
tb("전해질 (Electrolyte)", P3_X + 0.15, SCH_Y + 0.08,
   SCH_W3 - 0.20, 0.42, size=8, bold=True,
   color=rgb(0x01, 0x57, 0x9B))

# Cu 집전체
rect(P3_X + 0.05, CU_TOP - 0.85, SCH_W3, 0.85, fill=C_CU)
tb("Cu 집전체", P3_X + 0.05, CU_TOP - 0.65,
   SCH_W3, 0.45, size=9, bold=True,
   color=C_WHITE, align=PP_ALIGN.CENTER)

# ALD 불화물 인공 SEI (Cu 위 얇은 층)
ALD3_Y = CU_TOP - 0.85 - 0.28
rect(P3_X + 0.05, ALD3_Y, SCH_W3, 0.28, fill=C_ALD_FILM)
tb("ALD 불화물 인공 SEI  (ZnF₂ / AlF₃ / LiF)",
   P3_X + 0.05, ALD3_Y + 0.02, SCH_W3, 0.28,
   size=7.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 균일 Li 핵생성 (작은 반구형 원)
NUCL_Y = ALD3_Y - 0.32
NUCL_D = 0.30
for i, nx in enumerate([P3_X + 0.30 + i * 0.56 for i in range(14)]):
    circle(nx, NUCL_Y, NUCL_D, fill=C_LI)

# 평탄한 Li 증착층
LI3_Y = NUCL_Y - 1.30
rect(P3_X + 0.05, LI3_Y, SCH_W3, 1.30, fill=C_LI_FLAT)
tb("균일하고 평탄한 Li 증착층",
   P3_X + 0.05, LI3_Y + 0.42, SCH_W3, 0.45,
   size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 라벨
tb("✓ 균일 Li 핵생성", P3_X + 4.0, NUCL_Y - 0.40, 3.8, 0.45,
   size=9, bold=True, color=C_SOLUTION)
tb("인공 SEI (ALD 불화물)", P3_X + 4.0, ALD3_Y - 0.42, 3.8, 0.38,
   size=8, color=rgb(0x00, 0x69, 0x5C))

# § 설명 텍스트
DESC3_Y = SCH_Y + SCH_H + 0.10
rect(P3_X + 0.05, DESC3_Y, SCH_W3, 2.35,
     fill=rgb(0xE8, 0xF5, 0xE9), line=C_SOLUTION, lw=Pt(0.5))
mtb([
    "• ALD 불화물 인공 SEI 형성 후 Li 전착",
    "• 균일 핵생성 → 조밀·평탄한 Li 증착",
    "• 덴드라이트 억제, SEI 두께·조성 정밀 제어",
    "• 목표: CE >99.5%, 용량유지율 >80% (300 cy.)",
], P3_X + 0.15, DESC3_Y + 0.10, SCH_W3 - 0.25,
   color=rgb(0x1B, 0x5E, 0x20), size=8.5, lsp=13.0)

# § SEM 이미지 플레이스홀더
PH3_Y  = DESC3_Y + 2.50
PH3_H  = PANEL_H - (PH3_Y - PANEL_Y) - 0.12
rect(P3_X + 0.05, PH3_Y, SCH_W3, PH3_H,
     fill=C_PLACEHOLDER, line=C_GRID, lw=Pt(0.75))
line(P3_X + 0.05, PH3_Y, P3_X + 0.05 + SCH_W3, PH3_Y + PH3_H,
     color=C_GRID, width=Pt(0.5))
line(P3_X + 0.05 + SCH_W3, PH3_Y, P3_X + 0.05, PH3_Y + PH3_H,
     color=C_GRID, width=Pt(0.5))
tb("SEM 이미지 공간\n(균일 Li 증착 / 평탄 표면)",
   P3_X + 0.05, PH3_Y + PH3_H / 2 - 0.35, SCH_W3, 0.80,
   size=8, color=C_PH_TEXT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# § 6. 하단 보더
# ═══════════════════════════════════════════════════════════════
rect(0, H - 0.08, W, 0.08, fill=C_HEADER)


# ═══════════════════════════════════════════════════════════════
# § 7. 저장
# ═══════════════════════════════════════════════════════════════
OUT = (r"d:\23_GitHub\01_Office_PC\030_Project\032_Proposed Project"
       r"\09_개인기초 2차\그림1_연구개념도.pptx")
prs.save(OUT)
print(f"저장 완료: {OUT}")
