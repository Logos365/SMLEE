"""
ALD ZnF₂ 공정 사이클 모식도 (3D 스타일)
참조: ALD IAO 사이클 논문 그림 스타일
4-Step: DEZ펄스 → N₂퍼지 → HF·py펄스 → N₂퍼지(ZnF₂완성)
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import math

def rgb(r, g, b): return RGBColor(r, g, b)

# ── 색상 정의 ───────────────────────────────────────────────
C_ZN        = rgb(0xFF, 0xA0, 0x26)   # Zn: 주황색
C_ZN_DK     = rgb(0xCC, 0x66, 0x00)
C_F         = rgb(0x66, 0xBB, 0x6A)   # F: 초록
C_F_DK      = rgb(0x2E, 0x7D, 0x32)
C_ET        = rgb(0x78, 0x90, 0x9C)   # 에틸기 C: 청회색
C_ET_DK     = rgb(0x45, 0x62, 0x6F)
C_H         = rgb(0xF5, 0xF5, 0xDC)   # H: 베이지
C_PY        = rgb(0xCE, 0x93, 0xD8)   # Pyridine: 라벤더
C_N2        = rgb(0x90, 0xCA, 0xF9)   # N₂: 연파랑
C_SUB_TOP   = rgb(0x75, 0x75, 0x75)   # 기판 윗면: 중간회색
C_SUB_SIDE  = rgb(0x42, 0x42, 0x42)   # 기판 옆면: 어두운회색
C_SUB_BOT   = rgb(0x30, 0x30, 0x30)
C_BG        = rgb(0xFF, 0xFF, 0xFF)
C_ARROW     = rgb(0x42, 0x87, 0xC8)   # 사이클 화살표: 파랑
C_ARROW_DK  = rgb(0x1A, 0x57, 0xA0)
C_TITLE     = rgb(0x1A, 0x23, 0x7E)
C_LBL       = rgb(0x21, 0x21, 0x21)
C_STEP_BG1  = rgb(0xFF, 0xF3, 0xE0)   # DEZ step bg
C_STEP_BG2  = rgb(0xF3, 0xF3, 0xFF)   # purge step bg
C_STEP_BG3  = rgb(0xE8, 0xF5, 0xE9)   # HF step bg
C_STEP_BG4  = rgb(0xE3, 0xF2, 0xFD)   # complete step bg
C_BYPRODUCT = rgb(0xEF, 0x9A, 0x9A)   # 부산물 C₂H₆: 연적

# ── 슬라이드 설정 ──────────────────────────────────────────
W, H = 25.4, 19.05
prs = Presentation()
prs.slide_width  = Cm(W)
prs.slide_height = Cm(H)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# ── 기본 드로잉 함수 ───────────────────────────────────────
def shape(sid, x, y, w, h, fill, line=None, lw=Pt(0.5)):
    s = slide.shapes.add_shape(sid, Cm(x), Cm(y), Cm(w), Cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = lw
    else: s.line.fill.background()
    return s

def rect(x, y, w, h, fill, line=None, lw=Pt(0.5)):
    return shape(1, x, y, w, h, fill, line, lw)

def ellipse(x, y, w, h, fill, line=None, lw=Pt(0.5)):
    return shape(9, x, y, w, h, fill, line, lw)

def block_arrow(sid, x, y, w, h, fill, line=None):
    """Block arrow shapes"""
    return shape(sid, x, y, w, h, fill, line, Pt(0))

def tb(text, x, y, w, h, font="맑은 고딕", size=10, bold=False,
       color=C_LBL, align=PP_ALIGN.CENTER, wrap=True):
    txb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf  = txb.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size)
    run.font.bold = bold; run.font.color.rgb = color
    return txb

def line_connector(x1, y1, x2, y2, color=C_ET_DK, width=Pt(1.5)):
    conn = slide.shapes.add_connector(1, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    conn.line.color.rgb = color; conn.line.width = width
    return conn

# ── 3D 기판 그리기 ─────────────────────────────────────────
def draw_substrate(cx, sy, sw=9.2, eh=1.55, depth=0.55):
    """3D disk substrate  cx=center-x  sy=top-of-ellipse y"""
    ex = cx - sw / 2
    rect(ex, sy + eh * 0.45, sw, depth, fill=C_SUB_SIDE)
    ellipse(ex, sy + depth * 0.9, sw, eh, fill=C_SUB_BOT,
            line=C_SUB_BOT, lw=Pt(0.3))
    ellipse(ex, sy, sw, eh, fill=C_SUB_TOP,
            line=C_SUB_SIDE, lw=Pt(0.8))

# ── 원자 구 그리기 ─────────────────────────────────────────
def zn_ball(cx, cy, d=0.46):
    ellipse(cx-d/2, cy-d/2, d, d, fill=C_ZN,    line=C_ZN_DK,  lw=Pt(0.6))

def f_ball(cx, cy, d=0.32):
    ellipse(cx-d/2, cy-d/2, d, d, fill=C_F,     line=C_F_DK,   lw=Pt(0.5))

def c_ball(cx, cy, d=0.25):
    ellipse(cx-d/2, cy-d/2, d, d, fill=C_ET,    line=C_ET_DK,  lw=Pt(0.4))

def h_ball(cx, cy, d=0.18):
    ellipse(cx-d/2, cy-d/2, d, d, fill=C_H,     line=C_N2,     lw=Pt(0.3))

def py_ball(cx, cy, d=0.28):
    ellipse(cx-d/2, cy-d/2, d, d, fill=C_PY,    line=rgb(0x8E,0x24,0xAA), lw=Pt(0.4))

def n2_ball(cx, cy, d=0.26):
    ellipse(cx-d/2, cy-d/2, d, d, fill=C_N2,    line=rgb(0x19,0x76,0xD2), lw=Pt(0.4))

def et_ball(cx, cy, d=0.20):
    ellipse(cx-d/2, cy-d/2, d, d, fill=C_BYPRODUCT, line=rgb(0xC6,0x28,0x28), lw=Pt(0.3))

# ── DEZ 분자 (Zn(C₂H₅)₂) ─────────────────────────────────
def draw_dez(cx, cy, sc=1.0):
    """Diethylzinc: central Zn + 2 ethyl arms (simplified)"""
    d_zn = 0.44 * sc
    d_c  = 0.22 * sc
    arm  = 0.36 * sc
    # Left arm: C-C stick
    line_connector(cx - d_zn*0.4, cy, cx - d_zn*0.4 - arm, cy - arm*0.7,
                   color=C_ET_DK, width=Pt(1.5*sc))
    c_ball(cx - d_zn*0.4 - arm, cy - arm*0.7, d_c)
    line_connector(cx - d_zn*0.4 - arm, cy - arm*0.7,
                   cx - d_zn*0.4 - arm*2, cy - arm*1.2,
                   color=C_ET_DK, width=Pt(1.2*sc))
    c_ball(cx - d_zn*0.4 - arm*2.0, cy - arm*1.2, d_c * 0.85)
    # Right arm: C-C stick
    line_connector(cx + d_zn*0.4, cy, cx + d_zn*0.4 + arm, cy - arm*0.7,
                   color=C_ET_DK, width=Pt(1.5*sc))
    c_ball(cx + d_zn*0.4 + arm, cy - arm*0.7, d_c)
    line_connector(cx + d_zn*0.4 + arm, cy - arm*0.7,
                   cx + d_zn*0.4 + arm*2, cy - arm*1.2,
                   color=C_ET_DK, width=Pt(1.2*sc))
    c_ball(cx + d_zn*0.4 + arm*2.0, cy - arm*1.2, d_c * 0.85)
    # Central Zn
    zn_ball(cx, cy, d_zn)

# ── HF·py 분자 ────────────────────────────────────────────
def draw_hf_py(cx, cy, sc=1.0):
    """HF·pyridine: F + H connected"""
    d_f = 0.30 * sc; d_h = 0.18 * sc; d_py = 0.26 * sc
    line_connector(cx - d_f*0.3, cy, cx + d_f + d_h*0.3, cy,
                   color=C_F_DK, width=Pt(1.3*sc))
    f_ball(cx, cy, d_f)
    h_ball(cx + d_f + d_h*0.1, cy, d_h)
    # pyridine on F
    line_connector(cx - d_f*0.1, cy, cx - d_f - d_py, cy,
                   color=rgb(0x8E,0x24,0xAA), width=Pt(1.0*sc))
    py_ball(cx - d_f - d_py + d_py*0.1, cy, d_py)

# ── N₂ 분자 ───────────────────────────────────────────────
def draw_n2(cx, cy, sc=1.0):
    d = 0.24 * sc
    line_connector(cx - d*0.3, cy, cx + d*1.4, cy,
                   color=rgb(0x19,0x76,0xD2), width=Pt(1.5*sc))
    n2_ball(cx, cy, d)
    n2_ball(cx + d*0.9, cy, d)

# ── 부산물 (C₂H₆) ─────────────────────────────────────────
def draw_byproduct(cx, cy, sc=1.0):
    d = 0.20 * sc
    line_connector(cx, cy, cx + d*1.5, cy,
                   color=rgb(0xC6,0x28,0x28), width=Pt(1.2*sc))
    et_ball(cx, cy, d)
    et_ball(cx + d*1.2, cy, d)

# ── Zn 표면층 (기판 위) ────────────────────────────────────
def draw_zn_surface_row(cx, sy_surf, n=10, d=0.44, sp=0.87):
    """Row of Zn atoms adsorbed on surface"""
    total = (n - 1) * sp + d
    sx = cx - total / 2
    for i in range(n):
        zn_ball(sx + i * sp + d/2, sy_surf - d*0.35, d)

def draw_f_surface_row(cx, sy_surf, n=10, d=0.30, sp=0.87):
    """Row of F atoms on top of Zn layer"""
    total = (n - 1) * sp + d
    sx = cx - total / 2 + (sp - d) / 2
    for i in range(n):
        f_ball(sx + i * sp + d/2, sy_surf - d*0.3, d)

# ── 사이클 화살표 ─────────────────────────────────────────
def draw_cycle_arrow(sid, x, y, w, h):
    """Large block arrow for cycle flow"""
    s = slide.shapes.add_shape(sid, Cm(x), Cm(y), Cm(w), Cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = C_ARROW
    s.line.color.rgb = C_ARROW_DK; s.line.width = Pt(0.5)
    # Make it semi-transparent
    from pptx.oxml.ns import qn
    spPr = s.element
    solidFill = spPr.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        alpha = etree.SubElement(solidFill.find(qn('a:srgbClr')), qn('a:alpha'))
        alpha.set('val', '75000')  # ~75% opacity
    return s

# ═══════════════════════════════════════════════════════════════
# 레이아웃 좌표 계산
# ═══════════════════════════════════════════════════════════════
# 패널 (10.5 × 8.0 cm)
# Q1(TL): x=0.3,  y=0.4
# Q2(TR): x=14.6, y=0.4
# Q3(BR): x=14.6, y=10.65
# Q4(BL): x=0.3,  y=10.65
# 수평 갭: x=10.8 ~ 14.6 (3.8 cm)
# 수직 갭: y=8.4  ~ 10.65 (2.25 cm)
# ───────────────────────────────────────────────────────────────
PW, PH = 10.5, 8.0
Q1X, Q1Y = 0.3,  0.4
Q2X, Q2Y = 14.6, 0.4
Q3X, Q3Y = 14.6, 10.65
Q4X, Q4Y = 0.3,  10.65
PCX = [Q1X+PW/2, Q2X+PW/2, Q3X+PW/2, Q4X+PW/2]
PCY = [Q1Y+PH/2, Q2Y+PH/2, Q3Y+PH/2, Q4Y+PH/2]

GH_MX = (Q1X + PW + Q2X) / 2   # 수평 갭 중심 x = 12.7
GV_MY = (Q1Y + PH + Q3Y) / 2   # 수직 갭 중심 y ≈ 9.52

# 기판 y좌표 (각 패널 하단 기준)
SUB_EH = 1.55   # 기판 상면 타원 높이
SUB_DPT = 0.55  # 기판 두께
SUB_SW = 9.2    # 기판 너비

def sub_top_y(qy):
    return qy + PH - SUB_EH - SUB_DPT - 0.15

# 표면 원자 y (기판 상면 중심보다 약간 위)
def surf_y(qy):
    return sub_top_y(qy) + SUB_EH * 0.3

# ══════════════════════════════════════════════════════════════
# § 배경
# ══════════════════════════════════════════════════════════════
rect(0, 0, W, H, fill=C_BG)

# ══════════════════════════════════════════════════════════════
# § 패널 배경
# ══════════════════════════════════════════════════════════════
for (px, py, bg) in [
    (Q1X, Q1Y, C_STEP_BG1),
    (Q2X, Q2Y, C_STEP_BG2),
    (Q3X, Q3Y, C_STEP_BG3),
    (Q4X, Q4Y, C_STEP_BG4),
]:
    rect(px, py, PW, PH, fill=bg, line=rgb(0xCC, 0xCC, 0xCC), lw=Pt(0.5))

# ══════════════════════════════════════════════════════════════
# § 중앙 텍스트 영역
# ══════════════════════════════════════════════════════════════
rect(GH_MX - 2.5, GV_MY - 1.35, 5.0, 2.7,
     fill=rgb(0xF8, 0xF8, 0xFF), line=rgb(0x9F, 0xA8, 0xDA), lw=Pt(1.5))
tb("ALD ZnF₂", GH_MX - 2.4, GV_MY - 1.25, 4.8, 0.85,
   size=15, bold=True, color=C_TITLE)
tb("공정 사이클", GH_MX - 2.4, GV_MY - 0.40, 4.8, 0.75,
   size=13, bold=True, color=C_TITLE)
tb("(ZnF₂ ALD Cycle)", GH_MX - 2.4, GV_MY + 0.35, 4.8, 0.50,
   size=9, color=rgb(0x5C, 0x6B, 0xC0))

# ══════════════════════════════════════════════════════════════
# § 사이클 화살표 (4방향 블록 화살표)
# MSO: 13=RIGHT_ARROW 14=LEFT_ARROW 25=UP_ARROW 26=DOWN_ARROW
# ══════════════════════════════════════════════════════════════
ARW = 1.6  # 화살표 폭
ARH = 1.2  # 화살표 높이

# 위 → (Q1→Q2, 오른쪽)
draw_cycle_arrow(13, GH_MX - ARW/2, Q1Y + 2.0, ARW, ARH)
# 오른 ↓ (Q2→Q3, 아래)
draw_cycle_arrow(26, Q2X + PW*0.5 - ARH/2, GV_MY - ARW/2, ARH, ARW)
# 아래 ← (Q3→Q4, 왼쪽)
draw_cycle_arrow(14, GH_MX - ARW/2, Q4Y + 2.0, ARW, ARH)
# 왼 ↑ (Q4→Q1, 위)
draw_cycle_arrow(25, Q1X + PW*0.5 - ARH/2, GV_MY - ARW/2, ARH, ARW)


# ══════════════════════════════════════════════════════════════
# § STEP 1 — DEZ 펄스  (Q1: 상단 왼쪽)
# ══════════════════════════════════════════════════════════════
S1CX = PCX[0]
S1SY = sub_top_y(Q1Y)
S1FY = surf_y(Q1Y)

# 패널 타이틀
rect(Q1X, Q1Y, PW, 0.80, fill=rgb(0xFF, 0x8F, 0x00))
tb("① DEZ 펄스  (전구체 노출)", Q1X, Q1Y + 0.05, PW, 0.72,
   size=10, bold=True, color=C_BG)

# 기판
draw_substrate(S1CX, S1SY)

# 표면 Zn 원자 (부분 흡착 - 희소)
ZN_SP = 0.87
for i, xi in enumerate([S1CX - 3.9 + i * ZN_SP for i in range(10)]):
    if i % 3 != 1:  # 간헐적으로 흡착됨
        zn_ball(xi, S1FY - 0.15, 0.44)

# 기체상 DEZ 분자들 (기판 위 부유)
DEZ_POS = [
    (S1CX - 3.2, Q1Y + 1.5),
    (S1CX - 1.0, Q1Y + 2.8),
    (S1CX + 1.5, Q1Y + 1.3),
    (S1CX + 3.5, Q1Y + 2.5),
    (S1CX + 0.2, Q1Y + 3.8),
]
for (dx, dy) in DEZ_POS:
    draw_dez(dx, dy, sc=0.95)

# 하강 화살표 (분자가 기판으로 향함)
for (dx, dy) in [(S1CX - 1.0, Q1Y + 2.8), (S1CX + 1.5, Q1Y + 1.3)]:
    conn = slide.shapes.add_connector(
        1, Cm(dx), Cm(dy + 0.3), Cm(dx), Cm(S1FY - 0.5))
    conn.line.color.rgb = rgb(0xCC, 0x66, 0x00)
    conn.line.width = Pt(1.2)
    tail_xml = ('<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
                ' type="arrow" w="sm" len="sm"/>')
    conn.line._ln.append(etree.fromstring(tail_xml))

# 서브 레이블
tb("Zn(C₂H₅)₂ 전구체 흡착 개시",
   Q1X + 0.2, Q1Y + PH - 2.9, PW - 0.4, 0.45,
   size=8, color=rgb(0xBF, 0x36, 0x0C))


# ══════════════════════════════════════════════════════════════
# § STEP 2 — N₂ 퍼지  (Q2: 상단 오른쪽)
# ══════════════════════════════════════════════════════════════
S2CX = PCX[1]
S2SY = sub_top_y(Q2Y)
S2FY = surf_y(Q2Y)

rect(Q2X, Q2Y, PW, 0.80, fill=rgb(0x55, 0x6B, 0x9E))
tb("② N₂ 퍼지  (과량 DEZ 제거)", Q2X, Q2Y + 0.05, PW, 0.72,
   size=10, bold=True, color=C_BG)

draw_substrate(S2CX, S2SY)

# 표면 Zn (포화 흡착 - 촘촘)
for xi in [S2CX - 3.9 + i * ZN_SP for i in range(10)]:
    zn_ball(xi, S2FY - 0.15, 0.44)

# Zn에 남은 에틸 리간드 (절반)
for xi in [S2CX - 3.9 + i * ZN_SP for i in range(0, 10, 2)]:
    c_ball(xi - 0.3, S2FY - 0.55, 0.20)
    c_ball(xi - 0.5, S2FY - 0.82, 0.18)

# N₂ 화살표 (퍼지 흐름 표시)
for ny in [Q2Y + 1.8, Q2Y + 2.6, Q2Y + 3.4]:
    tb("N₂  →", Q2X + 0.2, ny, 3.0, 0.45,
       font="Arial", size=10, bold=True, color=C_N2, align=PP_ALIGN.LEFT)

# 제거되는 과잉 DEZ (흐릿하게 표시)
for (dx, dy) in [(S2CX - 2.5, Q2Y + 1.2), (S2CX + 2.0, Q2Y + 1.8)]:
    draw_dez(dx, dy, sc=0.7)
    # 제거 화살표 (위쪽으로)
    conn = slide.shapes.add_connector(
        1, Cm(dx), Cm(dy - 0.4), Cm(dx + 1.5), Cm(Q2Y + 0.9))
    conn.line.color.rgb = rgb(0x90, 0x90, 0x90)
    conn.line.width = Pt(1.0)
    tail_xml = ('<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
                ' type="arrow" w="sm" len="sm"/>')
    conn.line._ln.append(etree.fromstring(tail_xml))

tb("과잉 DEZ 제거 → 자기제한 표면",
   Q2X + 0.2, Q2Y + PH - 2.9, PW - 0.4, 0.45,
   size=8, color=rgb(0x1A, 0x23, 0x7E))


# ══════════════════════════════════════════════════════════════
# § STEP 3 — HF·py 펄스  (Q3: 하단 오른쪽)
# ══════════════════════════════════════════════════════════════
S3CX = PCX[2]
S3SY = sub_top_y(Q3Y)
S3FY = surf_y(Q3Y)

rect(Q3X, Q3Y, PW, 0.80, fill=rgb(0x2E, 0x7D, 0x32))
tb("③ HF·py 펄스  (반응물 노출)", Q3X, Q3Y + 0.05, PW, 0.72,
   size=10, bold=True, color=C_BG)

draw_substrate(S3CX, S3SY)

# 기판 위: Zn 흡착층
for xi in [S3CX - 3.9 + i * ZN_SP for i in range(10)]:
    zn_ball(xi, S3FY - 0.15, 0.44)

# HF와 반응 중 (일부 F 이미 결합)
for i, xi in enumerate([S3CX - 3.9 + i * ZN_SP for i in range(10)]):
    if i % 2 == 0:
        f_ball(xi, S3FY - 0.60, 0.30)  # F 이미 결합

# 기체상 HF·py 분자들
HF_POS = [
    (S3CX - 3.5, Q3Y + 1.6),
    (S3CX - 1.5, Q3Y + 2.9),
    (S3CX + 0.5, Q3Y + 1.5),
    (S3CX + 2.5, Q3Y + 2.7),
    (S3CX + 3.8, Q3Y + 1.8),
]
for (hx, hy) in HF_POS:
    draw_hf_py(hx, hy, sc=0.95)

# 하강 화살표
for (hx, hy) in [(S3CX - 1.5, Q3Y + 2.9), (S3CX + 0.5, Q3Y + 1.5)]:
    conn = slide.shapes.add_connector(
        1, Cm(hx), Cm(hy + 0.2), Cm(hx), Cm(S3FY - 0.5))
    conn.line.color.rgb = C_F_DK; conn.line.width = Pt(1.2)
    tail_xml = ('<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
                ' type="arrow" w="sm" len="sm"/>')
    conn.line._ln.append(etree.fromstring(tail_xml))

# 부산물 C₂H₆ 발생
for (bx, by) in [(S3CX - 2.0, Q3Y + 3.8), (S3CX + 1.8, Q3Y + 3.5)]:
    draw_byproduct(bx, by, sc=0.85)

tb("HF + Zn-C₂H₅* → ZnF* + C₂H₆↑",
   Q3X + 0.2, Q3Y + PH - 2.9, PW - 0.4, 0.45,
   size=8, color=rgb(0x1B, 0x5E, 0x20))


# ══════════════════════════════════════════════════════════════
# § STEP 4 — N₂ 퍼지 / ZnF₂ 단층 완성  (Q4: 하단 왼쪽)
# ══════════════════════════════════════════════════════════════
S4CX = PCX[3]
S4SY = sub_top_y(Q4Y)
S4FY = surf_y(Q4Y)

rect(Q4X, Q4Y, PW, 0.80, fill=rgb(0x15, 0x65, 0xC0))
tb("④ N₂ 퍼지 → ZnF₂ 단층 완성", Q4X, Q4Y + 0.05, PW, 0.72,
   size=10, bold=True, color=C_BG)

draw_substrate(S4CX, S4SY)

# ZnF₂ 완성 단층: Zn 행 + F 행 (2층)
for xi in [S4CX - 3.9 + i * ZN_SP for i in range(10)]:
    zn_ball(xi, S4FY - 0.15, 0.44)
for xi in [S4CX - 3.9 + i * ZN_SP for i in range(10)]:
    f_ball(xi, S4FY - 0.60, 0.30)

# 2번째 부분 단층 (다음 사이클 시작 예시)
for i, xi in enumerate([S4CX - 3.9 + i * ZN_SP for i in range(10)]):
    if i % 3 == 0:
        zn_ball(xi, S4FY - 1.00, 0.40)

# N₂ 퍼지 흐름
for ny in [Q4Y + 2.2, Q4Y + 3.0]:
    tb("N₂  →", Q4X + 0.2, ny, 3.2, 0.45,
       font="Arial", size=10, bold=True, color=C_N2, align=PP_ALIGN.LEFT)

# 완성 표시
rect(Q4X + PW*0.55, Q4Y + 1.5, PW*0.4, 0.55,
     fill=rgb(0xE3, 0xF2, 0xFD), line=rgb(0x15, 0x65, 0xC0), lw=Pt(1.0))
tb("✓ 1 사이클\n완료", Q4X + PW*0.55, Q4Y + 1.5, PW*0.4, 0.55,
   size=8, bold=True, color=rgb(0x0D, 0x47, 0xA1))

tb("ZnF₂ 단층 형성 → 다음 사이클로",
   Q4X + 0.2, Q4Y + PH - 2.9, PW - 0.4, 0.45,
   size=8, color=rgb(0x0D, 0x47, 0xA1))


# ══════════════════════════════════════════════════════════════
# § 범례 (하단 중앙)
# ══════════════════════════════════════════════════════════════
LEG_Y = H - 0.85
LEG_CX = W / 2
ITEMS = [
    (C_ZN,  C_ZN_DK,  "Zn",    0.34),
    (C_F,   C_F_DK,   "F",     0.28),
    (C_ET,  C_ET_DK,  "C (에틸기)", 0.24),
    (C_H,   C_N2,     "H",     0.20),
    (C_PY,  rgb(0x8E,0x24,0xAA), "Py (피리딘)", 0.24),
    (C_N2,  rgb(0x19,0x76,0xD2), "N₂",  0.22),
    (C_BYPRODUCT, rgb(0xC6,0x28,0x28), "C₂H₆ (부산물)", 0.20),
]
n_items = len(ITEMS)
item_w = W / (n_items + 1)
for i, (fc, lc, label, d) in enumerate(ITEMS):
    ix = item_w * (i + 0.8)
    ellipse(ix, LEG_Y, d, d, fill=fc, line=lc, lw=Pt(0.5))
    tb(label, ix - 0.1, LEG_Y + 0.05 + d, item_w * 0.9, 0.45,
       size=7.5, color=C_LBL)

# ══════════════════════════════════════════════════════════════
# § 저장
# ══════════════════════════════════════════════════════════════
OUT = (r"d:\23_GitHub\01_Office_PC\030_Project\032_Proposed Project"
       r"\09_개인기초 2차\그림_ALD_공정사이클.pptx")
prs.save(OUT)
print(f"저장 완료: {OUT}")
