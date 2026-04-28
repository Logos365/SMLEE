"""
그림 2 - 연구 로드맵 및 Task 구성도 PPTX 생성 스크립트
단일 슬라이드, A4 가로 비율 (25.4 cm × 17.78 cm)
"""

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── 색상 팔레트 ──────────────────────────────────────────────
C_HEADER     = RGBColor(0x28, 0x35, 0x93)   # 최상단 타이틀 바 (deep indigo)
C_TASK1_D    = RGBColor(0x15, 0x65, 0xC0)   # Task1 진파랑
C_TASK1_L    = RGBColor(0xE3, 0xF2, 0xFD)   # Task1 연파랑
C_TASK2_D    = RGBColor(0x2E, 0x7D, 0x32)   # Task2 진초록
C_TASK2_L    = RGBColor(0xE8, 0xF5, 0xE9)   # Task2 연초록
C_TASK3_D    = RGBColor(0xE6, 0x51, 0x00)   # Task3 진주황
C_TASK3_L    = RGBColor(0xFF, 0xF3, 0xE0)   # Task3 연주황
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BG         = RGBColor(0xF1, 0xF3, 0xF9)   # 상단 섹션 배경
C_DIVIDER    = RGBColor(0x39, 0x49, 0xAB)   # 구분선
C_GANTT_BG1  = RGBColor(0xF5, 0xF9, 0xFF)
C_GANTT_BG2  = RGBColor(0xF5, 0xFF, 0xF5)
C_GANTT_BG3  = RGBColor(0xFF, 0xFB, 0xF5)
C_YR1_HDR    = RGBColor(0xC5, 0xCA, 0xE9)
C_YR2_HDR    = RGBColor(0xBB, 0xDE, 0xFB)
C_YR3_HDR    = RGBColor(0xB3, 0xE5, 0xFC)
C_Q_HDR      = RGBColor(0xED, 0xE7, 0xF6)
C_GRID_MAJOR = RGBColor(0x79, 0x86, 0xCB)
C_GRID_MINOR = RGBColor(0xC5, 0xCA, 0xE9)
C_TASK1_MID  = RGBColor(0x19, 0x76, 0xD2)
C_TASK1_DK   = RGBColor(0x28, 0x35, 0x93)
C_TASK2_MID  = RGBColor(0x38, 0x8E, 0x3C)
C_TASK2_DK   = RGBColor(0x1B, 0x5E, 0x20)
C_TASK3_MID  = RGBColor(0xF5, 0x7C, 0x00)
C_TASK3_DK   = RGBColor(0xBF, 0x36, 0x0C)
C_ARROW      = RGBColor(0x45, 0x5A, 0x64)
C_LABEL      = RGBColor(0x3F, 0x51, 0xB5)
C_ROW_LBL1   = RGBColor(0xE3, 0xF2, 0xFD)
C_ROW_LBL2   = RGBColor(0xE8, 0xF5, 0xE9)
C_ROW_LBL3   = RGBColor(0xFF, 0xF3, 0xE0)
C_TARGET_BG  = RGBColor(0xFF, 0xE0, 0xB2)


def rgb(r, g, b):
    return RGBColor(r, g, b)


# ── 슬라이드 크기 ─────────────────────────────────────────────
W_CM = 25.4   # A4 가로
H_CM = 17.78  # A4 세로 (16:9 비율보다 약간 짧음)


# ── 헬퍼 함수 ─────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color=None, line_color=None,
             line_width=Pt(0), transparency=0):
    """단순 직사각형 추가"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Cm(x), Cm(y), Cm(w), Cm(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_name="맑은 고딕", font_size=10,
                 bold=False, color=RGBColor(0, 0, 0),
                 align=PP_ALIGN.LEFT, v_anchor=None,
                 word_wrap=True):
    """텍스트박스 추가"""
    txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_label_box(slide, text, x, y, w, h,
                  fill_color=None, text_color=C_WHITE,
                  font_size=9, bold=True, align=PP_ALIGN.CENTER,
                  line_color=None, line_width=Pt(0.5)):
    """배경색 있는 라벨 박스 (배지 용도)"""
    rect = add_rect(slide, x, y, w, h, fill_color, line_color, line_width)
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    return rect, tb


def add_multiline_textbox(slide, lines, x, y, w,
                          font_name="맑은 고딕", font_size=9,
                          color=RGBColor(0x1A, 0x23, 0x7E),
                          line_spacing_pt=13):
    """여러 줄 텍스트박스"""
    # 높이는 줄 수에 맞게 자동 계산
    h = len(lines) * (line_spacing_pt / 28.35)  # pt→cm rough conversion
    h = max(h, 0.5)
    txBox = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h + 0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        # 줄 간격
        from pptx.util import Pt as UPt
        p.line_spacing = UPt(line_spacing_pt)
    return txBox


# ── Presentation 생성 ─────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Cm(W_CM)
prs.slide_height = Cm(H_CM)

blank_layout = prs.slide_layouts[6]   # 완전 빈 레이아웃
slide = prs.slides.add_slide(blank_layout)

# ═══════════════════════════════════════════════════════════════
# 좌표계 (cm 단위)
# 전체 슬라이드: 0..25.4 × 0..17.78
# 상단 섹션 (Task 흐름): 0..0.95 → y=0~8.0 cm
# 하단 섹션 (간트):       y=8.0~17.78 cm
# ═══════════════════════════════════════════════════════════════

TOP_SEC_H = 8.0       # 상단 섹션 높이
HDR_H     = 1.15      # 최상단 타이틀 바

# 간트 좌표
GANTT_Y      = TOP_SEC_H + 0.1    # 간트 섹션 시작
GANTT_LBL_W  = 5.15               # 왼쪽 라벨 폭
GANTT_X      = GANTT_LBL_W        # 차트 영역 시작 x
GANTT_W      = W_CM - GANTT_LBL_W - 0.1   # 차트 영역 폭 (≈ 20.15cm)
GANTT_QW     = GANTT_W / 12       # 분기당 폭 (≈ 1.679cm)
GANTT_YR_H   = 0.66               # 연도 헤더 높이
GANTT_QH     = 0.45               # 분기 서브헤더 높이
GANTT_ROW_H  = 0.96               # 행 높이
GANTT_BAR_H  = 0.52               # 간트 바 높이
GANTT_BAR_DY = (GANTT_ROW_H - GANTT_BAR_H) / 2  # 바 수직 여백

# 분기 x좌표 함수
def qx(q):
    """0-based quarter index → x position (cm)"""
    return GANTT_X + q * GANTT_QW

ROWS_START_Y = GANTT_Y + GANTT_YR_H + GANTT_QH


# ═══════════════════════════════════════════════════════════════
# § 1. 최상단 타이틀 바
# ═══════════════════════════════════════════════════════════════
add_rect(slide, 0, 0, W_CM, HDR_H, fill_color=C_HEADER)
add_text_box(slide,
             "[그림 2]  연구 로드맵 및 Task 구성도",
             0, 0.05, W_CM, HDR_H - 0.1,
             font_name="맑은 고딕", font_size=13, bold=True,
             color=C_WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# § 2. 상단 섹션 배경 + 섹션 레이블
# ═══════════════════════════════════════════════════════════════
add_rect(slide, 0, HDR_H, W_CM, TOP_SEC_H - HDR_H, fill_color=C_BG)
add_text_box(slide, "▶ 연구 단계별 흐름",
             0.3, HDR_H + 0.05, 5, 0.4,
             font_name="맑은 고딕", font_size=8.5, bold=True,
             color=C_LABEL)


# ═══════════════════════════════════════════════════════════════
# § 3. Task 박스 공통 파라미터
# ═══════════════════════════════════════════════════════════════
BOX_Y  = HDR_H + 0.45   # 박스 시작 y
BOX_H  = TOP_SEC_H - BOX_Y - 0.15   # 박스 높이 ≈ 6.35 cm
BOX_W  = 7.2            # 박스 폭 (3개 → 7.2*3+gap*2=21.6+3.4=25, 약간 조정)
BOX_W  = 7.5
GAP    = (W_CM - 3 * BOX_W) / 4   # 좌우 여백 + 사이 간격 배분
BOX_X  = [GAP, GAP * 2 + BOX_W, GAP * 3 + BOX_W * 2]

HDR_H2     = 1.1    # 박스 헤더 높이
BADGE_H    = 0.42   # 연차 배지 높이
METHOD_H   = 0.40   # 분석법 배지 높이
BULLET_X   = 0.22   # 박스 내 텍스트 왼쪽 여백

def make_task_box(slide, bx, title1, title2, year_txt,
                  bullets, methods,
                  hdr_color, lbl_color, bg_color,
                  badge_color, method_color):
    """Task 박스 전체 생성"""
    # 외곽 박스 (배경)
    add_rect(slide, bx, BOX_Y, BOX_W, BOX_H,
             fill_color=bg_color, line_color=hdr_color, line_width=Pt(1.5))
    # 헤더 배경
    add_rect(slide, bx, BOX_Y, BOX_W, HDR_H2, fill_color=hdr_color)

    # 헤더 텍스트 line1
    add_text_box(slide, title1,
                 bx, BOX_Y + 0.05, BOX_W, 0.5,
                 font_name="맑은 고딕", font_size=10, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)
    # 헤더 텍스트 line2
    add_text_box(slide, title2,
                 bx, BOX_Y + 0.52, BOX_W, 0.5,
                 font_name="맑은 고딕", font_size=10, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)

    # 연차 배지
    badge_x = bx + (BOX_W - 3.0) / 2
    add_rect(slide, badge_x, BOX_Y + HDR_H2 + 0.1,
             3.0, BADGE_H,
             fill_color=badge_color)
    add_text_box(slide, year_txt,
                 badge_x, BOX_Y + HDR_H2 + 0.1,
                 3.0, BADGE_H,
                 font_name="맑은 고딕", font_size=8.5, bold=True,
                 color=C_WHITE, align=PP_ALIGN.CENTER)

    # 불릿 텍스트
    bullet_y = BOX_Y + HDR_H2 + BADGE_H + 0.28
    add_multiline_textbox(slide, bullets,
                          bx + BULLET_X, bullet_y,
                          BOX_W - BULLET_X * 2,
                          font_name="맑은 고딕", font_size=8.5,
                          color=lbl_color, line_spacing_pt=13.5)

    # 분석법 배지 (하단)
    method_y = BOX_Y + BOX_H - METHOD_H - 0.12
    mx = bx + 0.15
    for m_text, m_w in methods:
        add_rect(slide, mx, method_y, m_w, METHOD_H,
                 fill_color=method_color)
        add_text_box(slide, m_text,
                     mx, method_y, m_w, METHOD_H,
                     font_name="Arial", font_size=7.5, bold=True,
                     color=C_WHITE, align=PP_ALIGN.CENTER)
        mx += m_w + 0.10


# Task 1
make_task_box(
    slide, BOX_X[0],
    "Task 1.  ALD 불화물 박막 합성",
    "및 계면 특성 분석",
    "■  1년차",
    [
        "• 전구체·반응물 최적화 (ZnF₂, AlF₃, LiF)",
        "• 증착 조건 탐색 (두께 5~100 nm)",
        "• XRD·XPS·TEM/EDS 결정구조 분석",
        "• ToF-SIMS 계면 결합 상태 분석",
        "• EIS · Li⁺ 이온전도도 평가",
        "• Pinhole-free 합성 조건 확립",
    ],
    [("XRD", 0.85), ("XPS", 0.85), ("TEM", 0.85), ("EIS", 0.75), ("ToF-SIMS", 1.4)],
    hdr_color=C_TASK1_D, lbl_color=rgb(0x1A, 0x23, 0x7E),
    bg_color=C_TASK1_L, badge_color=rgb(0x0D, 0x47, 0xA1),
    method_color=C_TASK1_D,
)

# Task 2
make_task_box(
    slide, BOX_X[1],
    "Task 2.  In-situ/Operando 분석",
    "Li 핵생성 메커니즘 규명",
    "■  1~2년차",
    [
        "• Cryo-TEM·AFM  Li 핵생성 실시간 관찰",
        "• Operando XRD  충·방전 중 Li 거동",
        "• in-situ SEM  SEI 진화 모니터링",
        "• DFT: 불화물/Li 계면 확산 배리어",
        "• 과전압·증착 균일도 데이터베이스 구축",
    ],
    [("Cryo-TEM", 1.55), ("AFM", 0.85), ("Op.XRD", 1.25), ("DFT", 0.85), ("in-situ SEM", 1.7)],
    hdr_color=C_TASK2_D, lbl_color=rgb(0x1B, 0x5E, 0x20),
    bg_color=C_TASK2_L, badge_color=rgb(0x1B, 0x5E, 0x20),
    method_color=C_TASK2_D,
)

# Task 3
make_task_box(
    slide, BOX_X[2],
    "Task 3.  Full Cell 평가 및",
    "인공 SEI 설계 원칙 도출",
    "■  2~3년차",
    [
        "• Li||NCM 풀 셀 제작 및 장기 사이클",
        "  (>300 사이클,  CE >99.5%)",
        "• Post-mortem XPS/TEM SEI 분석",
        "• ZnF₂ vs AlF₃ vs LiF 비교 분석",
        "• ALD 기반 인공 SEI 설계 원칙 제안",
        "  목표: 용량유지율 >80% (300 사이클)",
    ],
    [("XPS", 0.85), ("TEM", 0.85), ("Full-Cell", 1.55), ("Post-mortem", 2.0)],
    hdr_color=C_TASK3_D, lbl_color=rgb(0xBF, 0x36, 0x0C),
    bg_color=C_TASK3_L, badge_color=rgb(0xBF, 0x36, 0x0C),
    method_color=C_TASK3_D,
)


# ═══════════════════════════════════════════════════════════════
# § 4. Task 간 연결 화살표
# ═══════════════════════════════════════════════════════════════
def add_arrow(slide, x1, y1, x2, y2, color=C_ARROW, width=Pt(2.5)):
    """직선 화살표"""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Cm(x1), Cm(y1), Cm(x2), Cm(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # 끝쪽 화살촉
    from pptx.oxml.ns import qn
    spPr = connector.line._ln
    # pptx XML로 화살촉 추가
    tail_xml = (
        '<a:tailEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' type="none"/>'
    )
    head_xml = (
        '<a:headEnd xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' type="arrow" w="med" len="med"/>'
    )
    spPr.append(etree.fromstring(tail_xml))
    spPr.append(etree.fromstring(head_xml))


ARR_Y = BOX_Y + BOX_H * 0.50
# 화살표 1: Task1 끝 → Task2 시작
add_arrow(slide, BOX_X[0] + BOX_W, ARR_Y,
                 BOX_X[1], ARR_Y)
# 화살표 레이블 1
mid1_x = (BOX_X[0] + BOX_W + BOX_X[1]) / 2
add_text_box(slide, "계면\n정보",
             mid1_x - 0.6, ARR_Y - 0.55, 1.2, 0.8,
             font_name="맑은 고딕", font_size=7.5,
             color=rgb(0x54, 0x6E, 0x7A), align=PP_ALIGN.CENTER)

# 화살표 2: Task2 끝 → Task3 시작
add_arrow(slide, BOX_X[1] + BOX_W, ARR_Y,
                 BOX_X[2], ARR_Y)
mid2_x = (BOX_X[1] + BOX_W + BOX_X[2]) / 2
add_text_box(slide, "메커니즘\n규명",
             mid2_x - 0.65, ARR_Y - 0.55, 1.3, 0.8,
             font_name="맑은 고딕", font_size=7.5,
             color=rgb(0x54, 0x6E, 0x7A), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# § 5. 구분선
# ═══════════════════════════════════════════════════════════════
add_rect(slide, 0, TOP_SEC_H, W_CM, 0.10, fill_color=C_DIVIDER)


# ═══════════════════════════════════════════════════════════════
# § 6. 간트 차트 섹션
# ═══════════════════════════════════════════════════════════════

# 섹션 레이블
add_text_box(slide, "▶ 연차별 연구 추진 일정 (간트 차트)",
             0.3, GANTT_Y + 0.05, 8, 0.42,
             font_name="맑은 고딕", font_size=8.5, bold=True,
             color=C_LABEL)

# ── 연도 헤더 ────────────────────────────────────────────────
YR_HDR_Y = GANTT_Y + 0.52
for i, (yr_txt, yr_clr, txt_clr) in enumerate([
    ("1년차", C_YR1_HDR, rgb(0x1A, 0x23, 0x7E)),
    ("2년차", C_YR2_HDR, rgb(0x0D, 0x47, 0xA1)),
    ("3년차", C_YR3_HDR, rgb(0x01, 0x57, 0x9B)),
]):
    yr_x = qx(i * 4)
    yr_w = GANTT_QW * 4
    add_rect(slide, yr_x, YR_HDR_Y, yr_w, GANTT_YR_H, fill_color=yr_clr)
    add_text_box(slide, yr_txt,
                 yr_x, YR_HDR_Y, yr_w, GANTT_YR_H,
                 font_name="맑은 고딕", font_size=10, bold=True,
                 color=txt_clr, align=PP_ALIGN.CENTER)

# ── 분기 서브헤더 ─────────────────────────────────────────────
Q_HDR_Y = YR_HDR_Y + GANTT_YR_H
add_rect(slide, GANTT_X, Q_HDR_Y, GANTT_W, GANTT_QH, fill_color=C_Q_HDR)
for qi in range(12):
    q_num = (qi % 4) + 1
    qc_x  = qx(qi)
    add_text_box(slide, f"{q_num}Q",
                 qc_x, Q_HDR_Y, GANTT_QW, GANTT_QH,
                 font_name="Arial", font_size=7.5, bold=False,
                 color=rgb(0x5C, 0x6B, 0xC0), align=PP_ALIGN.CENTER)

# ── 수직 격자선 ───────────────────────────────────────────────
for qi in range(13):
    lx = qx(qi)
    is_major = (qi % 4 == 0)
    lw = Pt(1.0) if is_major else Pt(0.5)
    lc = C_GRID_MAJOR if is_major else C_GRID_MINOR
    line_shape = slide.shapes.add_connector(
        1, Cm(lx), Cm(Q_HDR_Y), Cm(lx),
        Cm(ROWS_START_Y + GANTT_ROW_H * 9)
    )
    line_shape.line.color.rgb = lc
    line_shape.line.width = lw

# ── 행 정의 ───────────────────────────────────────────────────
# (라벨1, 라벨2, 간트바 시작Q(0-idx), 종료Q(0-idx, inclusive), 배경색, 바색)
rows = [
    # Task 1
    ("① ALD 증착 조건 최적화",         "(ZnF₂, AlF₃, LiF)",           0, 3,  C_ROW_LBL1, C_TASK1_D,   "전구체 최적화 및 ALD 파라미터 탐색"),
    ("② 박막 구조·화학 조성 분석",     "(XRD / XPS / TEM / ToF-SIMS)", 1, 5,  C_ROW_LBL1, C_TASK1_MID, "XRD / XPS / TEM-EDS / ToF-SIMS"),
    ("③ EIS · Li⁺ 이온전도도 평가",    "(pinhole-free 조건 확립)",       2, 3,  C_ROW_LBL1, C_TASK1_DK,  "EIS · Li⁺ 이온전도도"),
    # Task 2
    ("④ Cryo-TEM / AFM 핵생성 관찰",   "Li 핵생성 나노구조 실시간",      2, 7,  C_ROW_LBL2, C_TASK2_D,   "Cryo-TEM / AFM — Li 핵생성 나노구조 관찰"),
    ("⑤ Operando XRD / in-situ SEM",   "Li 증착·용해 거동 모니터링",     3, 7,  C_ROW_LBL2, C_TASK2_MID, "Operando XRD / in-situ SEM 모니터링"),
    ("⑥ DFT 이론 계산",                "계면 Li⁺ 확산 에너지 배리어",    3, 7,  C_ROW_LBL2, C_TASK2_DK,  "DFT: 불화물/Li 계면 Li⁺ 확산 에너지"),
    # Task 3
    ("⑦ Li||NCM 풀 셀 장기 사이클",    ">300 사이클, CE >99.5%",         6, 11, C_ROW_LBL3, C_TASK3_D,   "Li||NCM Full Cell 장기 사이클링 (>300 cycles)"),
    ("⑧ Post-mortem XPS/TEM 분석",     "SEI 조성 변화·열화 메커니즘",    8, 11, C_ROW_LBL3, C_TASK3_MID, "Post-mortem XPS / TEM 분석"),
    ("⑨ ALD 인공 SEI 설계 원칙 도출",  "ZnF₂/AlF₃/LiF 비교 정립",       9, 11, C_ROW_LBL3, C_TASK3_DK,  "범용 ALD-SEI 설계 원칙 정립"),
]

TASK_GROUP_COLOR = [C_TASK1_D] * 3 + [C_TASK2_D] * 3 + [C_TASK3_D] * 3

for ri, (lbl1, lbl2, q_start, q_end, bg_clr, bar_clr, bar_txt) in enumerate(rows):
    ry = ROWS_START_Y + ri * GANTT_ROW_H

    # 라벨 영역 배경
    add_rect(slide, 0, ry, GANTT_LBL_W, GANTT_ROW_H, fill_color=bg_clr)
    # Task 그룹 accent bar (좌측 5mm)
    add_rect(slide, 0, ry, 0.18, GANTT_ROW_H, fill_color=TASK_GROUP_COLOR[ri])

    # 라벨 텍스트
    lbl_color = [C_TASK1_D, C_TASK2_D, C_TASK3_D][[0,0,0,1,1,1,2,2,2][ri]]
    add_text_box(slide, lbl1,
                 0.25, ry + 0.04, GANTT_LBL_W - 0.3, 0.45,
                 font_name="맑은 고딕", font_size=7.8, bold=True,
                 color=lbl_color)
    add_text_box(slide, lbl2,
                 0.25, ry + 0.46, GANTT_LBL_W - 0.3, 0.42,
                 font_name="맑은 고딕", font_size=7.2, bold=False,
                 color=lbl_color)

    # 차트 영역 배경
    add_rect(slide, GANTT_X, ry, GANTT_W, GANTT_ROW_H,
             fill_color=rgb(0xFA, 0xFA, 0xFF) if ri < 3
             else (rgb(0xF5, 0xFF, 0xF5) if ri < 6 else rgb(0xFF, 0xFD, 0xF8)))

    # 간트 바
    bar_x = qx(q_start) + 0.04
    bar_w = GANTT_QW * (q_end - q_start + 1) - 0.08
    bar_y = ry + GANTT_BAR_DY
    add_rect(slide, bar_x, bar_y, bar_w, GANTT_BAR_H, fill_color=bar_clr)
    # 바 내부 텍스트
    add_text_box(slide, bar_txt,
                 bar_x, bar_y, bar_w, GANTT_BAR_H,
                 font_name="맑은 고딕", font_size=7.2, bold=False,
                 color=C_WHITE, align=PP_ALIGN.CENTER)

    # 수평 구분선
    is_group_border = (ri in [2, 5])
    lw = Pt(1.0) if is_group_border else Pt(0.4)
    lc = C_GRID_MAJOR if is_group_border else C_GRID_MINOR
    sep = slide.shapes.add_connector(
        1, Cm(0), Cm(ry + GANTT_ROW_H),
        Cm(W_CM), Cm(ry + GANTT_ROW_H)
    )
    sep.line.color.rgb = lc
    sep.line.width = lw

# ── 라벨/차트 세로 구분선 ────────────────────────────────────
vsep = slide.shapes.add_connector(
    1, Cm(GANTT_LBL_W), Cm(ROWS_START_Y),
    Cm(GANTT_LBL_W), Cm(ROWS_START_Y + GANTT_ROW_H * 9)
)
vsep.line.color.rgb = C_GRID_MAJOR
vsep.line.width = Pt(1.0)

# ── 최하단 보더 ──────────────────────────────────────────────
add_rect(slide, 0, H_CM - 0.08, W_CM, 0.08, fill_color=C_HEADER)


# ═══════════════════════════════════════════════════════════════
# § 7. 저장
# ═══════════════════════════════════════════════════════════════
OUT_PATH = r"d:\23_GitHub\01_Office_PC\030_Project\032_Proposed Project\09_개인기초 2차\그림2_연구로드맵.pptx"
prs.save(OUT_PATH)
print(f"저장 완료: {OUT_PATH}")
