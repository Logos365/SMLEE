#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
배치식 ALD vs 공압식(연속식) ALD 비교 PPT 생성
심사위원 대상 발표 자료
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ─────────────────────────────
# 색상 팔레트
# ─────────────────────────────
C_DARK_NAVY  = RGBColor(0x1B, 0x3A, 0x6B)   # 짙은 네이비 (배경/헤더)
C_BATCH_BLUE = RGBColor(0x1F, 0x6F, 0xBB)   # 배치식 = 파랑
C_BATCH_LIGHT= RGBColor(0xD6, 0xE8, 0xF9)   # 배치식 연한 배경
C_PNEU_ORANGE= RGBColor(0xD4, 0x5F, 0x0E)   # 공압식 = 주황
C_PNEU_LIGHT = RGBColor(0xFD, 0xE8, 0xD5)   # 공압식 연한 배경
C_GREEN      = RGBColor(0x21, 0x8C, 0x4F)   # 장점 녹색
C_GREEN_LIGHT= RGBColor(0xD5, 0xF0, 0xE0)
C_RED        = RGBColor(0xC0, 0x39, 0x2B)   # 단점 빨강
C_RED_LIGHT  = RGBColor(0xFA, 0xE0, 0xDD)
C_GRAY_BG    = RGBColor(0xF5, 0xF5, 0xF7)   # 슬라이드 배경
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY_TEXT  = RGBColor(0x55, 0x55, 0x55)
C_YELLOW     = RGBColor(0xFF, 0xC1, 0x07)
C_HEADER_BG  = RGBColor(0x1B, 0x3A, 0x6B)

FONT_KR  = "맑은 고딕"
FONT_NUM = "Arial"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─────────────────────────────
# 헬퍼 함수
# ─────────────────────────────

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # 완전 빈 레이아웃
    return prs.slides.add_slide(blank_layout)

def rect(slide, l, t, w, h, fill=None, line_color=None, lw_pt=0, radius=False):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shp = slide.shapes.add_shape(
        1,  # RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(lw_pt)
    else:
        shp.line.fill.background()
    return shp

def rounded_rect(slide, l, t, w, h, fill=None, line_color=None, lw_pt=0, adj=0.08):
    shp = slide.shapes.add_shape(
        5,  # ROUNDED RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    # 모서리 조정
    try:
        sp = shp._element
        prstGeom = sp.find(qn('p:spPr')).find(qn('a:prstGeom'))
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            gd_list = avLst.findall(qn('a:gd'))
            if gd_list:
                gd_list[0].set('fmla', f'val {int(adj*100000)}')
    except Exception:
        pass
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(lw_pt)
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, l, t, w, h, size=14, bold=False, color=C_BLACK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True, font=FONT_KR, line_spacing=None):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if line_spacing:
        from pptx.util import Pt as PT
        from pptx.oxml.ns import qn as QN
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return txBox

def txt_lines(slide, lines, l, t, w, h, size=12, bold=False, color=C_BLACK,
              align=PP_ALIGN.LEFT, font=FONT_KR, spacing=1.15):
    """여러 줄 텍스트박스"""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        is_bold = bold
        line_text = line
        if line.startswith("**") and line.endswith("**"):
            is_bold = True
            line_text = line[2:-2]
        run.text = line_text
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = color
        run.font.name = font
        # 줄간격
        from pptx.oxml.ns import qn as QN
        from lxml import etree as ET
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = ET.SubElement(pPr, QN('a:lnSpc'))
        spcPct = ET.SubElement(lnSpc, QN('a:spcPct'))
        spcPct.set('val', str(int(spacing * 100000)))
    return txBox

def arrow_down(slide, cx, y_top, height, color=C_GRAY_TEXT, width_pt=2):
    """세로 화살표 (아래 방향)"""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(cx), Inches(y_top),
        Inches(cx), Inches(y_top + height)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    return connector

def arrow_right(slide, x_left, cy, width, color=C_GRAY_TEXT, width_pt=2):
    connector = slide.shapes.add_connector(
        1,
        Inches(x_left), Inches(cy),
        Inches(x_left + width), Inches(cy)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)
    return connector

def slide_bg(slide, color=C_GRAY_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def header_bar(slide, title, subtitle=None, bar_color=C_HEADER_BG):
    rect(slide, 0, 0, 13.33, 1.15, fill=bar_color)
    txt(slide, title, 0.35, 0.15, 12.0, 0.55,
        size=24, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.72, 12.0, 0.38,
            size=13, bold=False, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.LEFT)

def side_label(slide, text, l, t, w, h, bg_color, text_color=C_WHITE, size=13, bold=True):
    rounded_rect(slide, l, t, w, h, fill=bg_color, adj=0.05)
    txt(slide, text, l, t, w, h, size=size, bold=bold, color=text_color, align=PP_ALIGN.CENTER)

def check_bullet(slide, items, l, t, w, size=12, color=C_BLACK, bullet="✔ ", gap=0.3, font=FONT_KR):
    for i, item in enumerate(items):
        txt(slide, bullet + item, l, t + i * gap, w, gap + 0.05,
            size=size, color=color, font=font)

def x_bullet(slide, items, l, t, w, size=12, color=C_RED, bullet="✖ ", gap=0.3, font=FONT_KR):
    for i, item in enumerate(items):
        txt(slide, bullet + item, l, t + i * gap, w, gap + 0.05,
            size=size, color=color, font=font)

# ═══════════════════════════════════════
# 슬라이드 1 — 표지
# ═══════════════════════════════════════
slide = blank_slide(prs)
slide_bg(slide, C_DARK_NAVY)

# 배경 장식 - 큰 도형
rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK_NAVY)
rect(slide, 0, 5.8, 13.33, 1.7, fill=RGBColor(0x12, 0x28, 0x52))

# 중앙 흰색 패널
rounded_rect(slide, 1.2, 1.0, 10.93, 4.5, fill=C_WHITE, adj=0.03)

# 상단 컬러 띠 (배치 / 공압)
rect(slide, 1.2, 1.0, 5.46, 0.5, fill=C_BATCH_BLUE)
rect(slide, 6.67, 1.0, 4.86, 0.5, fill=C_PNEU_ORANGE)

txt(slide, "배치식  ALD", 1.2, 1.0, 5.46, 0.5,
    size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txt(slide, "공압식(연속식)  ALD", 6.67, 1.0, 4.86, 0.5,
    size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# VS 원
rounded_rect(slide, 6.27, 0.85, 0.78, 0.78, fill=C_YELLOW, adj=0.5)
txt(slide, "VS", 6.27, 0.9, 0.78, 0.68,
    size=14, bold=True, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)

# 메인 제목
txt(slide, "배치식 ALD vs 공압식(연속식) ALD", 1.3, 1.7, 10.7, 0.85,
    size=28, bold=True, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)
txt(slide, "원리·구조·장단점 비교 — 심사위원을 위한 쉬운 설명", 1.3, 2.55, 10.7, 0.5,
    size=16, bold=False, color=C_GRAY_TEXT, align=PP_ALIGN.CENTER)

# 구분선
rect(slide, 3.0, 3.15, 7.33, 0.04, fill=RGBColor(0xCC, 0xCC, 0xCC))

# 키워드 3개
kw_data = [
    ("🔬  기술 원리", "자기제한 반응 기반 원자층 증착"),
    ("⚙️  공정 방식", "배치(정지) vs 연속(비행) 반응"),
    ("🏭  양산 전략", "레시피 개발 → 연속식 이전"),
]
for i, (kw, desc) in enumerate(kw_data):
    cx = 1.8 + i * 3.35
    rounded_rect(slide, cx, 3.3, 3.0, 1.6, fill=C_GRAY_BG, adj=0.05)
    txt(slide, kw, cx, 3.38, 3.0, 0.45,
        size=13, bold=True, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)
    txt(slide, desc, cx + 0.1, 3.82, 2.8, 0.7,
        size=11, color=C_GRAY_TEXT, align=PP_ALIGN.CENTER)

# 하단 서브 정보
txt(slide, "2026년 첨단소재 사업화 기술 개발(성과확장형)   |   2세부: 한국기계연구원 × 대성기계공업",
    0, 6.1, 13.33, 0.45,
    size=11, color=RGBColor(0x88, 0xAA, 0xDD), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# 슬라이드 2 — 목차
# ═══════════════════════════════════════
slide = blank_slide(prs)
slide_bg(slide, C_GRAY_BG)
header_bar(slide, "목  차", "배치식 ALD와 공압식 ALD, 어떻게 다를까요?")

toc_items = [
    ("01", "ALD란 무엇인가?",          "원자층 증착의 기본 원리",           C_BATCH_BLUE),
    ("02", "배치식 ALD",               "구조·원리·특징",                   C_BATCH_BLUE),
    ("03", "공압식(연속식) ALD",        "구조·원리·특징",                   C_PNEU_ORANGE),
    ("04", "한눈에 보는 비교",          "핵심 차이점 요약표",                C_DARK_NAVY),
    ("05", "장단점 상세 비교",          "각 방식의 강점과 한계",             C_GREEN),
    ("06", "본 과제에서의 역할 분담",   "기계연(배치식) × 대성기계(공압식)", C_PNEU_ORANGE),
]

for i, (num, title, desc, color) in enumerate(toc_items):
    row = i // 2
    col = i % 2
    lx = 0.5 + col * 6.45
    ty = 1.35 + row * 1.85

    rounded_rect(slide, lx, ty, 6.1, 1.6, fill=C_WHITE, adj=0.04)
    rect(slide, lx, ty, 0.55, 1.6, fill=color)
    txt(slide, num, lx, ty, 0.55, 1.6,
        size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER, font=FONT_NUM)
    txt(slide, title, lx + 0.65, ty + 0.15, 5.2, 0.5,
        size=15, bold=True, color=C_DARK_NAVY)
    txt(slide, desc, lx + 0.65, ty + 0.65, 5.2, 0.6,
        size=12, color=C_GRAY_TEXT)

# ═══════════════════════════════════════
# 슬라이드 3 — ALD 기본 원리
# ═══════════════════════════════════════
slide = blank_slide(prs)
slide_bg(slide, C_GRAY_BG)
header_bar(slide, "01   ALD란 무엇인가?", "Atomic Layer Deposition — 원자층 증착 기술")

# 정의 박스
rounded_rect(slide, 0.4, 1.25, 12.53, 1.05, fill=C_BATCH_LIGHT, adj=0.04)
rect(slide, 0.4, 1.25, 0.12, 1.05, fill=C_BATCH_BLUE)
txt(slide,
    "ALD는 기체 전구체(Precursor)를 한 종류씩 순서대로 공급하여, 표면에서 '자기제한 반응(Self-limiting)'으로 원자 한 층씩 정밀하게 박막을 쌓는 증착 기술입니다.",
    0.62, 1.3, 12.1, 0.95, size=13, color=C_DARK_NAVY)

# ALD 사이클 다이어그램
step_data = [
    ("Step 1\n전구체 A\n공급 (Pulse)", C_BATCH_BLUE,   "전구체 A 분자가\n표면에 흡착"),
    ("Step 2\n불활성 기체\n퍼지 (Purge)", C_GRAY_TEXT,  "미반응 전구체 A\n제거"),
    ("Step 3\n반응제 B\n공급 (H₂O)", C_PNEU_ORANGE,   "A와 B가 반응 →\n박막 1층 형성"),
    ("Step 4\n불활성 기체\n퍼지 (Purge)", C_GRAY_TEXT,  "부산물 및 잉여\n반응제 제거"),
]
bx = 0.5
for i, (label, color, desc) in enumerate(step_data):
    cx = bx + i * 3.05
    rounded_rect(slide, cx, 2.5, 2.7, 1.4, fill=color, adj=0.06)
    txt(slide, label, cx, 2.5, 2.7, 1.4,
        size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, desc, cx + 0.1, 4.0, 2.5, 0.7,
        size=11, color=C_GRAY_TEXT, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(slide, "→", cx + 2.72, 3.0, 0.35, 0.5,
            size=18, bold=True, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)

# 반복 표시
txt(slide, "⟳  1 사이클 반복 → 원자 1층 증착  (GPC = 0.1 ~ 2.0 Å/cycle)",
    0.5, 4.75, 12.33, 0.45,
    size=13, bold=True, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)

# 핵심 특징 3개
feat_data = [
    ("📐 Å 수준 정밀 제어", "1 사이클 = 정확히 0.1~2.0 Å\n두께를 사이클 수로 결정"),
    ("🔄 자기제한 반응", "과잉 전구체 자동 제거\n→ 항상 균일한 1층 형성"),
    ("🔗 완전 균일 피복", "복잡한 형상·분말 입자\n표면까지 100% 균일 코팅"),
]
for i, (title, desc) in enumerate(feat_data):
    fx = 0.5 + i * 4.25
    rounded_rect(slide, fx, 5.3, 4.0, 1.75, fill=C_WHITE, adj=0.05)
    rect(slide, fx, 5.3, 4.0, 0.42, fill=C_DARK_NAVY)
    txt(slide, title, fx, 5.3, 4.0, 0.42,
        size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, desc, fx + 0.15, 5.78, 3.7, 1.2,
        size=12, color=C_GRAY_TEXT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# 슬라이드 4 — 배치식 ALD
# ═══════════════════════════════════════
slide = blank_slide(prs)
slide_bg(slide, C_GRAY_BG)
header_bar(slide, "02   배치식 ALD (Batch ALD)", "분말이 반응기 안에서 '정지' 상태로 반응 — 한국기계연구원 보유 장비", bar_color=C_BATCH_BLUE)

# 왼쪽: 구조 다이어그램
rect(slide, 0.35, 1.3, 5.9, 5.85, fill=C_WHITE)
txt(slide, "배치식 ALD 반응기 구조", 0.35, 1.35, 5.9, 0.4,
    size=13, bold=True, color=C_BATCH_BLUE, align=PP_ALIGN.CENTER)

# 반응기 본체
rounded_rect(slide, 1.2, 1.95, 4.2, 3.1, fill=RGBColor(0xE8, 0xF2, 0xFC), adj=0.05,
             line_color=C_BATCH_BLUE, lw_pt=2)
txt(slide, "반응기 (Reactor)", 1.2, 1.95, 4.2, 0.42,
    size=12, bold=True, color=C_BATCH_BLUE, align=PP_ALIGN.CENTER)

# 분말 (정지 상태)
rounded_rect(slide, 1.8, 2.65, 2.9, 1.7, fill=RGBColor(0xB3, 0xD4, 0xF1), adj=0.1)
txt(slide, "● ● ● ● ●\n분말 입자\n(정지 상태)", 1.8, 2.7, 2.9, 1.6,
    size=12, bold=True, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)

# 가스 입출구
txt(slide, "↓ 전구체 A / 반응제 B / 퍼지 가스", 1.5, 2.2, 3.6, 0.4,
    size=10, color=C_GRAY_TEXT, align=PP_ALIGN.CENTER)
txt(slide, "↑ 배기 (진공펌프)", 1.5, 4.38, 3.6, 0.38,
    size=10, color=C_GRAY_TEXT, align=PP_ALIGN.CENTER)

# 공정 흐름
txt(slide, "공정 순서 (Time-Domain ALD)", 0.5, 5.15, 5.6, 0.4,
    size=12, bold=True, color=C_BATCH_BLUE)
steps = ["① 전구체 A Pulse → 분말 표면 흡착 (정지 상태)",
         "② 퍼지 → 미반응 A 제거",
         "③ 반응제 B(H₂O) Pulse → 분말 표면 반응",
         "④ 퍼지 → 부산물 제거  ⟳  반복"]
for i, s in enumerate(steps):
    bullet_color = C_BATCH_BLUE if i in [0, 2] else C_GRAY_TEXT
    txt(slide, s, 0.6, 5.58 + i * 0.37, 5.5, 0.38,
        size=11, color=bullet_color)

# 오른쪽: 특징
rect(slide, 6.45, 1.3, 6.5, 5.85, fill=RGBColor(0xF0, 0xF5, 0xFF))

txt(slide, "배치식 ALD 핵심 특징", 6.55, 1.38, 6.3, 0.42,
    size=14, bold=True, color=C_BATCH_BLUE)

char_items = [
    ("🔹 처리 방식", "분말이 반응기 내부에 정지\n전구체가 분말 주위를 채움"),
    ("🔹 두께 제어", "Pulse 시간 = 전구체 노출 시간\n→ 직접적이고 정밀한 제어 가능"),
    ("🔹 공정 개발", "레시피 개발이 매우 용이\n★★★★★ (최고 수준)"),
    ("🔹 처리 용량", "소량 ~ 중량 (g ~ kg 수준)\n주로 R&D 및 파일럿 단계"),
    ("🔹 운전 환경", "진공 조건 필요\n배치 단위로 시작/종료"),
    ("🔹 본 과제 역할", "TiO₂, Al₂O₃, V₂O₅, ZnO\n4종 ALD 레시피 개발 소스"),
]
for i, (label, content) in enumerate(char_items):
    ty = 1.9 + i * 0.85
    txt(slide, label, 6.55, ty, 2.0, 0.38, size=12, bold=True, color=C_BATCH_BLUE)
    txt(slide, content, 6.55, ty + 0.35, 6.2, 0.45, size=11, color=C_GRAY_TEXT)

# 강조 박스
rounded_rect(slide, 6.45, 6.53, 6.5, 0.55, fill=C_BATCH_LIGHT, adj=0.05)
txt(slide, "🏛  한국기계연구원(KIMM) — 배치식 ALD 장비 2대 보유 · 운용 중",
    6.55, 6.57, 6.3, 0.45,
    size=12, bold=True, color=C_BATCH_BLUE)

# ═══════════════════════════════════════
# 슬라이드 5 — 공압식 ALD
# ═══════════════════════════════════════
slide = blank_slide(prs)
slide_bg(slide, C_GRAY_BG)
header_bar(slide, "03   공압식(연속식) ALD (Pneumatic Continuous ALD)",
           "분말이 기류(공압)를 타고 '비행'하며 반응 — 대성기계 OPERIO 10", bar_color=C_PNEU_ORANGE)

# 왼쪽 구조 다이어그램
rect(slide, 0.35, 1.3, 5.9, 5.85, fill=C_WHITE)
txt(slide, "공압식 ALD 반응기 구조 (구간 분리형)", 0.35, 1.35, 5.9, 0.4,
    size=13, bold=True, color=C_PNEU_ORANGE, align=PP_ALIGN.CENTER)

# 이송 파이프 (수평)
zone_data = [
    ("전구체 A\n구간", C_BATCH_BLUE,   "분말 ↔ A 반응"),
    ("퍼지\n구간",    C_GRAY_TEXT,     "A 제거"),
    ("반응제 B\n구간", C_PNEU_ORANGE,  "분말 ↔ B 반응"),
    ("퍼지\n구간",    C_GRAY_TEXT,     "부산물 제거"),
]
py = 2.7
for i, (label, color, sub) in enumerate(zone_data):
    zx = 0.55 + i * 1.36
    rect(slide, zx, py, 1.25, 1.0, fill=color)
    txt(slide, label, zx, py, 1.25, 1.0,
        size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, sub, zx, py + 1.05, 1.25, 0.38,
        size=9, color=C_GRAY_TEXT, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(slide, "→", zx + 1.27, py + 0.3, 0.12, 0.4,
            size=14, bold=True, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)

txt(slide, "← 공압(Air Pressure) 이송 →", 0.55, 2.3, 5.45, 0.38,
    size=11, bold=True, color=C_PNEU_ORANGE, align=PP_ALIGN.CENTER)

# 분말 이동 표시
txt(slide, "분말 입자가 기류를 타고 각 반응 구간을 '비행'하며 통과", 0.5, 3.85, 5.5, 0.45,
    size=11, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)

# 스케일업 단계
txt(slide, "3단계 스케일업 로드맵", 0.5, 4.42, 5.6, 0.4,
    size=12, bold=True, color=C_PNEU_ORANGE)
scale_steps = [
    ("1차년도", "R&D FBR (소량)", "두께 편차 ≤10%"),
    ("2차년도", "OPERIO 10 Pilot", "두께 편차 ≤7%"),
    ("3차년도", "양산  10 kg/h", "두께 편차 ≤5%"),
]
for i, (year, device, target) in enumerate(scale_steps):
    sx = 0.55 + i * 1.88
    rounded_rect(slide, sx, 4.88, 1.75, 1.5, fill=C_PNEU_LIGHT, adj=0.05)
    txt(slide, year, sx, 4.92, 1.75, 0.38,
        size=11, bold=True, color=C_PNEU_ORANGE, align=PP_ALIGN.CENTER)
    txt(slide, device, sx, 5.32, 1.75, 0.5,
        size=10, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)
    txt(slide, target, sx, 5.82, 1.75, 0.5,
        size=10, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    if i < 2:
        txt(slide, "→", sx + 1.77, 5.45, 0.12, 0.4,
            size=14, bold=True, color=C_PNEU_ORANGE, align=PP_ALIGN.CENTER)

# 오른쪽 특징
rect(slide, 6.45, 1.3, 6.5, 5.85, fill=RGBColor(0xFF, 0xF5, 0xEF))
txt(slide, "공압식 ALD 핵심 특징", 6.55, 1.38, 6.3, 0.42,
    size=14, bold=True, color=C_PNEU_ORANGE)

char_items = [
    ("🔸 처리 방식", "분말이 공압(기류)을 타고 이동\n각 반응 구간을 '비행'하며 코팅"),
    ("🔸 두께 제어", "RTD(체류시간 분포) 기반 간접 제어\nPulse 시간 ≠ 노출 시간 → 변환 필요"),
    ("🔸 공정 개발", "레시피 개발 난이도 높음\n★★☆☆☆ (배치식 레시피 이전 필요)"),
    ("🔸 처리 용량", "연속 운전 → 10 kg/h 이상\n산업 양산 직접 적용 가능"),
    ("🔸 운전 환경", "상압(대기압) 연속 운전 가능\n24시간 연속 운전 실증 목표"),
    ("🔸 본 과제 역할", "기계연 레시피를 RTD 변환 후 적용\n10 kg/h 양산 플랫폼"),
]
for i, (label, content) in enumerate(char_items):
    ty = 1.9 + i * 0.85
    txt(slide, label, 6.55, ty, 2.0, 0.38, size=12, bold=True, color=C_PNEU_ORANGE)
    txt(slide, content, 6.55, ty + 0.35, 6.2, 0.45, size=11, color=C_GRAY_TEXT)

# 강조 박스
rounded_rect(slide, 6.45, 6.53, 6.5, 0.55, fill=C_PNEU_LIGHT, adj=0.05)
txt(slide, "🏭  대성기계공업 OPERIO 10 — 국내 유일 연속식 분말 ALD 장비 보유",
    6.55, 6.57, 6.3, 0.45,
    size=12, bold=True, color=C_PNEU_ORANGE)

# ═══════════════════════════════════════
# 슬라이드 6 — 한눈에 보는 비교표
# ═══════════════════════════════════════
slide = blank_slide(prs)
slide_bg(slide, C_GRAY_BG)
header_bar(slide, "04   한눈에 보는 비교", "배치식 vs 공압식(연속식) ALD 핵심 차이점")

# 헤더 컬럼
col_w = [3.8, 4.3, 4.3]
col_x = [0.35, 4.2, 8.55]
# 헤더
rect(slide, col_x[0], 1.35, col_w[0], 0.55, fill=C_DARK_NAVY)
txt(slide, "비교 항목", col_x[0], 1.35, col_w[0], 0.55,
    size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
rect(slide, col_x[1], 1.35, col_w[1], 0.55, fill=C_BATCH_BLUE)
txt(slide, "🔵  배치식 ALD", col_x[1], 1.35, col_w[1], 0.55,
    size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
rect(slide, col_x[2], 1.35, col_w[2], 0.55, fill=C_PNEU_ORANGE)
txt(slide, "🟠  공압식(연속식) ALD", col_x[2], 1.35, col_w[2], 0.55,
    size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 데이터 행
rows = [
    ("분말 상태",         "정지(Static)\n반응기 내 고정",               "비행(In-flight)\n기류에 의해 이송"),
    ("반응 방식",         "시간 분리형(Time-domain)\nA → 퍼지 → B → 퍼지",  "공간 분리형(Spatial)\n구간별 A·B·퍼지 분리"),
    ("두께 제어",         "직접 제어\nPulse 시간 = 노출 시간",           "RTD 기반 간접 제어\n(μ - 2σ ≥ t_pulse)"),
    ("레시피 개발 난이도", "★★★★★  매우 용이",                        "★★☆☆☆  어려움\n(배치 레시피 이전 필요)"),
    ("처리량(Throughput)", "소~중량 (g~kg)\nR&D 최적",                  "대량 (10 kg/h 이상)\n양산 직접 적용"),
    ("운전 환경",         "진공 조건 필요\n배치 단위 시작/종료",          "상압 연속 운전\n24시간 연속 가능"),
    ("본 과제 역할",       "레시피 소스 (기계연)\nTiO₂·Al₂O₃·V₂O₅·ZnO 4종", "양산 플랫폼 (대성기계)\n10 kg/h 실증"),
]
row_colors = [C_WHITE, C_GRAY_BG, C_WHITE, C_GRAY_BG, C_WHITE, C_GRAY_BG, C_BATCH_LIGHT]
row_h = 0.72

for ri, (item, batch, pneu) in enumerate(rows):
    ty = 1.9 + ri * row_h
    bg = row_colors[ri]
    for ci in range(3):
        rect(slide, col_x[ci], ty, col_w[ci], row_h, fill=bg,
             line_color=RGBColor(0xCC, 0xCC, 0xCC), lw_pt=0.5)

    txt(slide, item, col_x[0] + 0.1, ty + 0.05, col_w[0] - 0.15, row_h - 0.1,
        size=12, bold=True, color=C_DARK_NAVY)
    txt(slide, batch, col_x[1] + 0.1, ty + 0.05, col_w[1] - 0.15, row_h - 0.1,
        size=11, color=C_BATCH_BLUE if ri < 6 else C_DARK_NAVY)
    txt(slide, pneu, col_x[2] + 0.1, ty + 0.05, col_w[2] - 0.15, row_h - 0.1,
        size=11, color=C_PNEU_ORANGE if ri < 6 else C_DARK_NAVY)

# 마지막 행 강조 border
rect(slide, col_x[0], 1.9 + 6 * row_h, sum(col_w) + 0.45, row_h,
     fill=None, line_color=C_DARK_NAVY, lw_pt=2)

# ═══════════════════════════════════════
# 슬라이드 7 — 장단점 상세 비교
# ═══════════════════════════════════════
slide = blank_slide(prs)
slide_bg(slide, C_GRAY_BG)
header_bar(slide, "05   장단점 상세 비교", "각 방식의 강점과 한계 — 왜 두 방식을 함께 사용하는가?")

# 배치식 ALD 패널
rounded_rect(slide, 0.35, 1.35, 6.05, 5.85, fill=C_WHITE, adj=0.04)
rect(slide, 0.35, 1.35, 6.05, 0.5, fill=C_BATCH_BLUE)
txt(slide, "🔵  배치식 ALD", 0.35, 1.35, 6.05, 0.5,
    size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 배치 장점
rounded_rect(slide, 0.5, 1.98, 5.75, 0.38, fill=C_GREEN, adj=0.05)
txt(slide, "✅  장점 (Strengths)", 0.55, 1.98, 5.65, 0.38,
    size=12, bold=True, color=C_WHITE)
batch_pros = [
    "레시피 개발이 매우 용이하고 재현성 높음",
    "Å 수준 두께 정밀 제어 가능 (직접 측정)",
    "공정 조건 최적화 후 연속식으로 이전 가능",
    "다양한 소재(TiO₂·Al₂O₃·V₂O₅·ZnO) 레시피 병렬 개발",
    "HR-TEM/XPS 등 분석과 바로 연계 가능",
]
for i, s in enumerate(batch_pros):
    txt(slide, "  ◆  " + s, 0.5, 2.43 + i * 0.38, 5.75, 0.38,
        size=11, color=C_GREEN)

# 배치 단점
rounded_rect(slide, 0.5, 4.4, 5.75, 0.38, fill=C_RED, adj=0.05)
txt(slide, "❌  단점 (Limitations)", 0.55, 4.4, 5.65, 0.38,
    size=12, bold=True, color=C_WHITE)
batch_cons = [
    "처리량 제한 — 소·중량(g~kg), 양산 직접 적용 불가",
    "배치 단위 운전 → 대량 생산 시 비효율",
    "진공 장비 필요 → 초기 투자 비용 높음",
]
for i, s in enumerate(batch_cons):
    txt(slide, "  ▲  " + s, 0.5, 4.85 + i * 0.42, 5.75, 0.42,
        size=11, color=C_RED)

# 공압식 ALD 패널
rounded_rect(slide, 6.6, 1.35, 6.38, 5.85, fill=C_WHITE, adj=0.04)
rect(slide, 6.6, 1.35, 6.38, 0.5, fill=C_PNEU_ORANGE)
txt(slide, "🟠  공압식(연속식) ALD", 6.6, 1.35, 6.38, 0.5,
    size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

rounded_rect(slide, 6.75, 1.98, 6.08, 0.38, fill=C_GREEN, adj=0.05)
txt(slide, "✅  장점 (Strengths)", 6.8, 1.98, 5.9, 0.38,
    size=12, bold=True, color=C_WHITE)
pneu_pros = [
    "대량 연속 처리 — 10 kg/h 이상 양산 가능",
    "상압 연속 운전 → 24시간 연속 생산 실현",
    "공압 이송으로 분말 뭉침 방지, 균일 분산",
    "구간 분리형 구조 → 전구체 교차 오염 방지",
    "스케일업 용이 — R&D → Pilot → 양산 단계적 확장",
]
for i, s in enumerate(pneu_pros):
    txt(slide, "  ◆  " + s, 6.75, 2.43 + i * 0.38, 6.08, 0.38,
        size=11, color=C_GREEN)

rounded_rect(slide, 6.75, 4.4, 6.08, 0.38, fill=C_RED, adj=0.05)
txt(slide, "❌  단점 (Limitations)", 6.8, 4.4, 5.9, 0.38,
    size=12, bold=True, color=C_WHITE)
pneu_cons = [
    "레시피 개발 어려움 — 비행 중 노출시간 직접 측정 불가",
    "RTD 변환 알고리즘 필요 (배치 레시피 그대로 적용 불가)",
    "초기 두께 편차 관리 복잡 (RTD 분포 관리 필요)",
]
for i, s in enumerate(pneu_cons):
    txt(slide, "  ▲  " + s, 6.75, 4.85 + i * 0.42, 6.08, 0.42,
        size=11, color=C_RED)

# 하단 연계 설명
rounded_rect(slide, 0.35, 6.6, 12.63, 0.68, fill=C_DARK_NAVY, adj=0.04)
txt(slide,
    "💡  두 방식의 상호 보완:  기계연 배치식(레시피 개발)  →  RTD 기반 변환  →  대성기계 공압식(양산 적용)  — 이것이 본 과제의 핵심 파이프라인",
    0.5, 6.65, 12.3, 0.55,
    size=12, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# 슬라이드 8 — 역할 분담 & RTD 변환
# ═══════════════════════════════════════
slide = blank_slide(prs)
slide_bg(slide, C_GRAY_BG)
header_bar(slide, "06   본 과제에서의 역할 분담",
           "배치식(기계연) × 공압식(대성기계) — RTD 기반 레시피 이전 파이프라인")

# 기계연 박스
rounded_rect(slide, 0.35, 1.35, 5.5, 5.2, fill=C_BATCH_LIGHT, adj=0.04,
             line_color=C_BATCH_BLUE, lw_pt=2)
rect(slide, 0.35, 1.35, 5.5, 0.52, fill=C_BATCH_BLUE)
txt(slide, "🏛  한국기계연구원 (KIMM)", 0.35, 1.35, 5.5, 0.52,
    size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

kimm_items = [
    ("배치식 ALD 장비 2대 운용", "TiO₂·Al₂O₃·V₂O₅·ZnO 4종 레시피 개발"),
    ("Self-limiting 반응 확인", "GPC vs Pulse 시간 포화 곡선 측정"),
    ("독창 구조 최초 개발", "V₂O₅(내층)/TiO₂(외층) ALD 이중층"),
    ("종합 분석 인프라", "HR-TEM·EELS·XPS·Raman·FTIR"),
    ("레시피 이전 방법론 확립", "배치→연속식 RTD 정량 변환 알고리즘"),
]
for i, (title, desc) in enumerate(kimm_items):
    ty = 2.05 + i * 0.85
    rounded_rect(slide, 0.5, ty, 5.15, 0.78, fill=C_WHITE, adj=0.04)
    txt(slide, "▶ " + title, 0.65, ty + 0.04, 4.9, 0.35,
        size=12, bold=True, color=C_BATCH_BLUE)
    txt(slide, "    " + desc, 0.65, ty + 0.38, 4.9, 0.35,
        size=11, color=C_GRAY_TEXT)

# 화살표 (중간)
for ay in [2.4, 3.3, 4.2, 5.1, 5.95]:
    txt(slide, "→", 5.9, ay, 1.5, 0.38,
        size=18, bold=True, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)

# 중앙 RTD 변환 박스
rounded_rect(slide, 5.85, 1.35, 1.6, 5.2, fill=RGBColor(0xFD, 0xF6, 0xDC), adj=0.04,
             line_color=C_YELLOW, lw_pt=2)
txt(slide, "RTD\n변환", 5.85, 1.45, 1.6, 0.8,
    size=13, bold=True, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)
formula_lines = [
    "μ - 2σ",
    "≥",
    "t_pulse",
]
for i, fl in enumerate(formula_lines):
    txt(slide, fl, 5.9, 2.5 + i * 0.55, 1.5, 0.5,
        size=12, bold=True, color=C_DARK_NAVY, align=PP_ALIGN.CENTER, font=FONT_NUM)
txt(slide, "체류시간\n분포(RTD)\n기반 변환", 5.9, 4.3, 1.5, 1.2,
    size=10, color=C_GRAY_TEXT, align=PP_ALIGN.CENTER)

# 대성기계 박스
rounded_rect(slide, 7.5, 1.35, 5.45, 5.2, fill=C_PNEU_LIGHT, adj=0.04,
             line_color=C_PNEU_ORANGE, lw_pt=2)
rect(slide, 7.5, 1.35, 5.45, 0.52, fill=C_PNEU_ORANGE)
txt(slide, "🏭  대성기계공업 (OPERIO 10)", 7.5, 1.35, 5.45, 0.52,
    size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

daes_items = [
    ("공압 이송 연속식 ALD 운용", "분말이 기류 타고 비행하며 코팅"),
    ("RTD 기반 레시피 적용", "기계연 레시피 → 연속식 변환 적용"),
    ("구간 분리형 반응기", "전구체 A/퍼지/반응제 B 물리적 분리"),
    ("10 kg/h 양산 실증", "24시간 연속 운전, 두께 편차 ≤5%"),
    ("Turn-key 공급 역량", "레시피 + 장비 + SOP 패키지"),
]
for i, (title, desc) in enumerate(daes_items):
    ty = 2.05 + i * 0.85
    rounded_rect(slide, 7.65, ty, 5.15, 0.78, fill=C_WHITE, adj=0.04)
    txt(slide, "▶ " + title, 7.78, ty + 0.04, 4.9, 0.35,
        size=12, bold=True, color=C_PNEU_ORANGE)
    txt(slide, "    " + desc, 7.78, ty + 0.38, 4.9, 0.35,
        size=11, color=C_GRAY_TEXT)

# 하단 결론 박스
rounded_rect(slide, 0.35, 6.65, 12.63, 0.72, fill=C_DARK_NAVY, adj=0.04)
txt(slide,
    "✅ 기계연(배치식) = 레시피 개발·분석 전문   |   대성기계(공압식) = 양산 스케일업 전문   →   두 기관의 결합으로 세계 최초 배치→연속 이전 방법론 확립",
    0.5, 6.7, 12.3, 0.6,
    size=11, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# 슬라이드 9 — 핵심 기술 난제와 해결방안
# ═══════════════════════════════════════
slide = blank_slide(prs)
slide_bg(slide, C_GRAY_BG)
header_bar(slide, "핵심 기술 난제 &  해결 방안",
           "배치식 → 공압식 이전의 3대 기술 장벽과 본 컨소시엄의 해결책")

prob_data = [
    (
        "난제 ①\n배치→연속 공정 이전",
        "배치식: 분말 정지 → Pulse 시간 = 노출 시간\n연속식: 분말 비행 → Pulse 시간 ≠ 노출 시간\n→ 기존 레시피 그대로 적용 시 두께 편차 >30%",
        "RTD 정량 변환 알고리즘\n최소 체류시간(μ-2σ) ≥ t_pulse\n세계 최초 배치→연속 정량 이전 방법론",
        "기계연 + 대성기계",
        "GPC 재현성 ≤5%\n(3차년도)",
    ),
    (
        "난제 ②\nV₂O₅ 전기화학 활성",
        "V₂O₅ 작동전압(2.0~3.5V)이\nSIB 양극 운전전압과 중복\n→ 두꺼울수록 독립 전극으로 반응",
        "ALD 두께 1~2 nm 정밀 제어\n비정질(amorphous) VOₓ 선호\nV₂O₅/TiO₂ 이중층으로 TiO₂ 외층이 차단",
        "기계연",
        "계면저항 ≤50 Ω·cm²\n(3차년도)",
    ),
    (
        "난제 ③\n양산 균일도 확보",
        "g 수준 → 10 kg/h : 처리량 3,000배 증가\n분말 뭉침·정체구간 발생\n전구체 가스 불균일 분포",
        "공압 이송으로 항상 분산 유지\n구간 분리형 반응기 구조\nRTD 기반 실시간 드리프트 보정",
        "대성기계",
        "두께 편차 ≤5%\n@ 10 kg/h",
    ),
]

for ci, (prob_title, prob_desc, sol_desc, owner, kpi) in enumerate(prob_data):
    cx = 0.38 + ci * 4.32

    # 난제 박스
    rounded_rect(slide, cx, 1.38, 4.1, 2.3, fill=C_RED_LIGHT, adj=0.05,
                 line_color=C_RED, lw_pt=1.5)
    rect(slide, cx, 1.38, 4.1, 0.48, fill=C_RED)
    txt(slide, prob_title, cx, 1.38, 4.1, 0.48,
        size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, prob_desc, cx + 0.12, 1.93, 3.86, 1.7,
        size=10.5, color=C_BLACK, align=PP_ALIGN.LEFT)

    # 화살표
    txt(slide, "▼", cx + 1.6, 3.72, 0.9, 0.42,
        size=16, bold=True, color=C_DARK_NAVY, align=PP_ALIGN.CENTER)

    # 해결 박스
    rounded_rect(slide, cx, 4.18, 4.1, 2.25, fill=C_GREEN_LIGHT, adj=0.05,
                 line_color=C_GREEN, lw_pt=1.5)
    rect(slide, cx, 4.18, 4.1, 0.48, fill=C_GREEN)
    txt(slide, "✅ 해결 방안", cx, 4.18, 4.1, 0.48,
        size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, sol_desc, cx + 0.12, 4.72, 3.86, 1.2,
        size=10.5, color=C_BLACK, align=PP_ALIGN.LEFT)

    # 담당 / KPI
    txt(slide, "담당: " + owner, cx + 0.12, 5.98, 2.1, 0.35,
        size=10, bold=True, color=C_DARK_NAVY)
    txt(slide, "목표: " + kpi, cx + 2.1, 5.98, 2.0, 0.35,
        size=10, bold=True, color=C_GREEN)

# ═══════════════════════════════════════
# 슬라이드 10 — 결론 / 요약
# ═══════════════════════════════════════
slide = blank_slide(prs)
slide_bg(slide, C_DARK_NAVY)

# 배경 장식
rect(slide, 0, 0, 13.33, 7.5, fill=C_DARK_NAVY)
rect(slide, 0, 6.7, 13.33, 0.8, fill=RGBColor(0x12, 0x28, 0x52))

txt(slide, "결  론  및  요  약", 0, 0.35, 13.33, 0.65,
    size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
txt(slide, "배치식 ALD + 공압식 ALD — 상호 보완으로 완성하는 양산 파이프라인",
    0, 1.0, 13.33, 0.48,
    size=14, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)

# 3개 요약 카드
summary_data = [
    ("🔵  배치식 ALD\n(한국기계연구원)",
     "• 분말 정지 상태 → 정밀 레시피 개발\n• Å 수준 두께 직접 제어\n• 4종 산화물 레시피 체계 확립\n• HR-TEM/XPS 원자 단위 분석",
     C_BATCH_BLUE, C_BATCH_LIGHT),
    ("⟶  RTD 변환 브릿지  ⟶",
     "• 배치식 t_pulse\n  → 연속식 μ-2σ 매핑\n• 세계 최초\n  정량적 이전 방법론\n• 두께 편차 ≤5% 목표",
     C_DARK_NAVY, RGBColor(0xFD, 0xF6, 0xDC)),
    ("🟠  공압식 ALD\n(대성기계 OPERIO 10)",
     "• 분말 비행 이송 → 연속 코팅\n• 10 kg/h 양산 처리량\n• 24시간 연속 운전 실증\n• 두께 편차 ≤5% 달성",
     C_PNEU_ORANGE, C_PNEU_LIGHT),
]
for i, (title, content, hdr_color, bg_color) in enumerate(summary_data):
    sx = 0.5 + i * 4.3
    rounded_rect(slide, sx, 1.65, 4.0, 4.3, fill=bg_color, adj=0.04)
    rect(slide, sx, 1.65, 4.0, 0.65, fill=hdr_color)
    txt(slide, title, sx, 1.65, 4.0, 0.65,
        size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, content, sx + 0.15, 2.4, 3.7, 3.4,
        size=12, color=C_BLACK, align=PP_ALIGN.LEFT)

# KPI 요약
kpi_box_data = [
    ("ALD 코팅 균일도", "두께 편차 ≤5%", "10 kg/h 연속 운전 기준"),
    ("GPC 재현성", "≤5%", "24시간 연속 운전 기준"),
    ("계면저항", "≤50 Ω·cm²", "EIS 평가, 사이클 후"),
    ("초기 쿨롱 효율", "≥85%", "1Ah 파우치셀 기준"),
]
for i, (label, target, note) in enumerate(kpi_box_data):
    kx = 0.5 + i * 3.2
    rounded_rect(slide, kx, 6.1, 3.0, 0.95, fill=RGBColor(0x1F, 0x4B, 0x8C), adj=0.05)
    txt(slide, label, kx, 6.15, 3.0, 0.32,
        size=10, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)
    txt(slide, target, kx, 6.45, 3.0, 0.38,
        size=14, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
    txt(slide, note, kx, 6.83, 3.0, 0.25,
        size=9, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)

txt(slide, "2026년 첨단소재 사업화 기술 개발(성과확장형)  |  2세부: 기계연 × 대성기계",
    0, 7.2, 13.33, 0.28,
    size=10, color=RGBColor(0x55, 0x77, 0xAA), align=PP_ALIGN.CENTER)

# ─────────────────────────────
# 저장
# ─────────────────────────────
output_path = r"d:\23_GitHub\01_Office_PC\030_Project\032_Proposed Project\08_2026년 첨단소재 사업화 기술 개발(성과확장형)\ALD_배치식_vs_공압식_비교.pptx"
prs.save(output_path)
print(f"[완료] 저장됨: {output_path}")
print(f"슬라이드 수: {len(prs.slides)}")
