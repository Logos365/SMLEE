"""
ALD 기술난제 해결방안 PPT 생성 스크립트
python-pptx + matplotlib 사용
"""
import sys, io, math
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.oxml.ns import qn
from lxml import etree
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patches as FancyArrowPatch
import numpy as np
from matplotlib.patches import FancyArrow, FancyArrowPatch
import matplotlib.font_manager as fm
import warnings
warnings.filterwarnings('ignore')

# ─── 폰트 설정 (한글 지원) ────────────────────────────────────
def setup_korean_font():
    """윈도우 한글 폰트 설정"""
    font_candidates = [
        'Malgun Gothic', 'NanumGothic', 'AppleGothic', 'sans-serif'
    ]
    for font in font_candidates:
        if any(font.lower() in f.name.lower() for f in fm.fontManager.ttflist):
            plt.rcParams['font.family'] = font
            return font
    plt.rcParams['font.family'] = 'Malgun Gothic'
    return 'Malgun Gothic'

setup_korean_font()
plt.rcParams['axes.unicode_minus'] = False

# ─── 색상 팔레트 ────────────────────────────────────────────────
C_NAVY      = RGBColor(0x1F, 0x4E, 0x79)   # 진한 네이비
C_BLUE      = RGBColor(0x2E, 0x75, 0xB6)   # 미디엄 블루
C_LIGHTBLUE = RGBColor(0xBD, 0xD7, 0xEE)   # 연한 블루
C_RED       = RGBColor(0xC0, 0x00, 0x00)   # 진한 빨강
C_LIGHTRED  = RGBColor(0xFF, 0xCC, 0xCC)   # 연한 빨강
C_GREEN     = RGBColor(0x37, 0x56, 0x23)   # 진한 녹색
C_LIGHTGREEN= RGBColor(0xE2, 0xEF, 0xDA)   # 연한 녹색
C_ORANGE    = RGBColor(0xED, 0x7D, 0x31)   # 오렌지
C_YELLOW    = RGBColor(0xFF, 0xC0, 0x00)   # 노랑
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)   # 흰색
C_BLACK     = RGBColor(0x00, 0x00, 0x00)   # 검정
C_GRAY      = RGBColor(0x7F, 0x7F, 0x7F)   # 회색
C_LIGHTGRAY = RGBColor(0xF2, 0xF2, 0xF2)   # 연회색
C_DARKGRAY  = RGBColor(0x40, 0x40, 0x40)   # 어두운 회색
C_ACCENT    = RGBColor(0xFF, 0x80, 0x00)   # 강조 오렌지

# ─── 유틸 함수 ──────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb, line_rgb=None, line_w=0):
    """사각형 도형 추가"""
    shape = slide.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_w if line_w else 1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, l, t, w, h, text, font_size=12, bold=False,
                 color=C_BLACK, align=PP_ALIGN.LEFT, wrap=True,
                 v_anchor=None, italic=False):
    """텍스트 박스 추가"""
    txb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    # 여백 최소화
    tf.margin_left  = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top   = Pt(2)
    tf.margin_bottom= Pt(2)

    if '\n' in text:
        lines = text.split('\n')
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = lines[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        for line in lines[1:]:
            p2 = tf.add_paragraph()
            p2.alignment = align
            r2 = p2.add_run()
            r2.text = line
            r2.font.size = Pt(font_size)
            r2.font.bold = bold
            r2.font.italic = italic
            r2.font.color.rgb = color
    else:
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return txb

def add_label_box(slide, l, t, w, h, label, value,
                  lbl_size=9, val_size=14, bg_fill=C_LIGHTBLUE,
                  lbl_color=C_NAVY, val_color=C_NAVY, border=C_BLUE):
    """라벨 + 값 박스"""
    box = add_rect(slide, l, t, w, h, bg_fill, border, 1.5)
    add_text_box(slide, l+0.1, t+0.1, w-0.2, h*0.35,
                 label, lbl_size, False, lbl_color, PP_ALIGN.CENTER)
    add_text_box(slide, l+0.1, t+h*0.38, w-0.2, h*0.55,
                 value, val_size, True,  val_color, PP_ALIGN.CENTER)
    return box

def add_arrow_down(slide, cx, top_y, bot_y, color=C_NAVY, width_cm=0.4):
    """아래 방향 화살표"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    arrow = slide.shapes.add_shape(
        13,  # DOWN_ARROW
        Cm(cx - width_cm/2), Cm(top_y),
        Cm(width_cm), Cm(bot_y - top_y)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow

def fig_to_image(fig):
    """matplotlib figure → bytes"""
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor='none', transparent=True)
    buf.seek(0)
    plt.close(fig)
    return buf

def add_fig(slide, fig, l, t, w, h):
    """figure를 슬라이드에 삽입"""
    img_buf = fig_to_image(fig)
    slide.shapes.add_picture(img_buf, Cm(l), Cm(t), Cm(w), Cm(h))

# ─── 공통 슬라이드 배경/제목 함수 ────────────────────────────────
def setup_slide_background(slide, prs):
    """흰 배경 + 상단 네이비 헤더 바"""
    SW = prs.slide_width  / 914400 * 2.54  # cm
    SH = prs.slide_height / 914400 * 2.54
    # 흰 배경
    add_rect(slide, 0, 0, SW, SH, C_WHITE)
    # 상단 네이비 헤더 바
    add_rect(slide, 0, 0, SW, 1.6, C_NAVY)
    # 하단 연한 파란 줄
    add_rect(slide, 0, SH - 0.5, SW, 0.5, C_LIGHTBLUE)
    return SW, SH

def add_slide_title(slide, prs, title, subtitle=None):
    """슬라이드 헤더에 제목 추가"""
    SW = prs.slide_width / 914400 * 2.54
    add_text_box(slide, 0.4, 0.15, SW - 0.8, 0.9,
                 title, 20, True, C_WHITE, PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, 0.4, 0.95, SW - 0.8, 0.5,
                     subtitle, 10, False, C_LIGHTBLUE, PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════
#  슬라이드 생성
# ═══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)   # 와이드 16:9
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

SW_TOTAL = prs.slide_width  / 914400 * 2.54   # ≈ 33.87 cm
SH_TOTAL = prs.slide_height / 914400 * 2.54   # ≈ 19.05 cm


# ─────────────────────────────────────────────────────────────────
# SLIDE 1 — 표지 슬라이드
# ─────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)

# 배경 그라데이션 (네이비→블루 느낌: 2개 레이어)
add_rect(s1, 0, 0, SW_TOTAL, SH_TOTAL, C_NAVY)
add_rect(s1, 0, SH_TOTAL*0.55, SW_TOTAL, SH_TOTAL*0.45, C_BLUE)

# 중앙 흰색 구분선
add_rect(s1, 2, SH_TOTAL*0.48, SW_TOTAL-4, 0.06, C_YELLOW)

# 제목
add_text_box(s1, 2, 3.5, SW_TOTAL-4, 3,
             "연속식 분말 ALD 코팅의\n3대 기술적 난제와 해결 방안",
             36, True, C_WHITE, PP_ALIGN.CENTER)

# 부제
add_text_box(s1, 2, 8.8, SW_TOTAL-4, 1.2,
             "한국기계연구원 (KIMM)  ×  대성기계공업㈜",
             18, False, C_LIGHTBLUE, PP_ALIGN.CENTER)

# 배경 키워드 (우측 하단 장식)
add_text_box(s1, SW_TOTAL-14, SH_TOTAL-3.5, 13.5, 3,
             "Batch-to-Continuous Recipe Transfer\n"
             "V₂O₅/TiO₂ Bilayer ALD  ·  RTD Algorithm\n"
             "10 kg/h Continuous Powder ALD",
             11, False, RGBColor(0x60, 0x90, 0xC0), PP_ALIGN.RIGHT)

# 과제명 하단
add_rect(s1, 0, SH_TOTAL-1.1, SW_TOTAL, 1.1, RGBColor(0x0D, 0x27, 0x47))
add_text_box(s1, 0.5, SH_TOTAL-0.95, SW_TOTAL-1, 0.85,
             "2026년 첨단소재 원천기술 성장지원사업 (2026-나노소재-6차)  |  2세부 선정평가 발표자료",
             9, False, C_LIGHTBLUE, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────
# SLIDE 2 — 3대 난제 개요 (Overview)
# ─────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
SW, SH = setup_slide_background(s2, prs)
add_slide_title(s2, prs,
                "ALD 코팅 기술 — 왜 양산이 어려운가?",
                "산업화의 3대 핵심 장벽 및 본 과제의 해결 접근")

# 배경 설명 박스
add_rect(s2, 0.4, 1.85, SW-0.8, 1.1, C_LIGHTBLUE, C_BLUE, 1)
add_text_box(s2, 0.6, 1.95, SW-1.2, 0.85,
    "ALD(원자층 증착)는 Å 수준 두께 제어와 완전 균일 피복(conformal coating)으로 SIB 양극재 계면 안정화에 "
    "이상적인 기술이다.\n"
    "그러나 배치식→양산 이전 불가 · 코팅 활성 제어 곤란 · 대량 균일 코팅 부재 — 3가지 난제가 상용화의 핵심 장벽이다.",
    10, False, C_NAVY, PP_ALIGN.LEFT)

# 3개 난제 박스
col_w = (SW - 1.2) / 3
col_xs = [0.4 + i * col_w for i in range(3)]
tops = [3.2, 3.2, 3.2]
titles = ["난제 ①", "난제 ②", "난제 ③"]
subtitles = [
    "배치식 → 연속식\n공정 이전 불가",
    "V₂O₅ 코팅층의\n전기화학적 활성",
    "대량 양산 시\n코팅 균일도 미확보"
]
icons = ["⚙", "⚡", "📦"]
descs = [
    "배치식(정지반응) vs 연속식(비행반응)\n공정 환경의 근본적 차이로\n레시피 직접 이전 불가능\n→ 두께 편차 >30% 발생",
    "V₂O₅ 작동전압(2.0–3.5V)이\nSIB 운전전압(2.0–4.2V)과 중복\n→ 두꺼운 코팅 시 성능 저하\n→ Å 수준 두께 제어 필수",
    "분말 뭉침(Agglomeration)\n정체구간(Dead Zone) 발생\n전구체 노출 불균일\n→ g수준→10 kg/h: 3,000배↑"
]
for i in range(3):
    x = col_xs[i] + 0.1
    bw = col_w - 0.2

    # 헤더 빨간 박스
    add_rect(s2, x, tops[i], bw, 1.0, C_RED)
    add_text_box(s2, x, tops[i]+0.05, bw, 0.45,
                 titles[i], 16, True, C_WHITE, PP_ALIGN.CENTER)
    add_text_box(s2, x, tops[i]+0.48, bw, 0.48,
                 subtitles[i], 10, True, RGBColor(0xFF,0xEE,0xEE), PP_ALIGN.CENTER)

    # 설명 연한 박스
    add_rect(s2, x, tops[i]+1.05, bw, 3.2, C_LIGHTRED, C_RED, 0.75)
    add_text_box(s2, x+0.15, tops[i]+1.2, bw-0.3, 3.0,
                 descs[i], 10.5, False, C_DARKGRAY, PP_ALIGN.LEFT)

# 하단 메시지
add_rect(s2, 0.4, SH-1.3, SW-0.8, 0.95, C_NAVY, C_NAVY, 0)
add_text_box(s2, 0.6, SH-1.25, SW-1.2, 0.8,
    "→  이 3가지 난제를 동시에 해결하기 위해, 본 과제는 기계연(ALD 레시피) × 대성기계(연속 장비)의 전문화된 협력 체계를 구성",
    10.5, True, C_YELLOW, PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────
# SLIDE 3 — 3열 메인 (난제 + 해결방안)
# ─────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
SW, SH = setup_slide_background(s3, prs)
add_slide_title(s3, prs,
                "3대 기술적 난제 및 해결 방안",
                "기계연(레시피·분석) × 대성기계(연속장비·양산) 협력 모델")

col_w = (SW - 1.2) / 3
col_xs = [0.4 + i * col_w for i in range(3)]

problem_data = [
    ("난제 ①\n배치 → 연속 이전 불가",
     "• 배치식: 분말 정지 → Pulse시간 = 노출시간\n"
     "• 연속식: 분말 비행 중 → Pulse시간 ≠ 노출시간\n"
     "• 직접 이전 시 두께 편차 >30%\n"
     "  → 재현성 확보 불가"),
    ("난제 ②\nV₂O₅ 전기화학적 활성",
     "• V₂O₅ 운전전압: 2.0–3.5V\n"
     "• SIB 양극 운전전압: 2.0–4.2V\n"
     "• 전압범위 중복 → 두꺼우면 독립 전극\n"
     "  → 전압 프로파일 교란 → 성능 저하"),
    ("난제 ③\n양산 시 균일도 미확보",
     "• 처리량: g 수준 → 10 kg/h (×3,000배)\n"
     "• 분말 뭉침 → 미코팅 입자 발생\n"
     "• 정체구간(Dead Zone) → 과코팅 혼재\n"
     "  → GPC 편차 증가 → 품질 저하"),
]
solution_data = [
    ("해결 ①\nRTD 기반 정량 변환 알고리즘",
     "• 체류시간 분포(RTD) 측정\n"
     "  최소 체류시간: μ - 2σ ≥ t_pulse\n"
     "• 배치식 변수 → 연속식 변수 1:1 매핑\n"
     "• 기계연 레시피 → 대성기계 이전\n"
     "★ 세계 최초 정량적 이전 방법론 확립"),
    ("해결 ②\nALD 이중층 설계로 활성 차단",
     "• ALD 자기제한: 1사이클 ≈ 0.5Å\n"
     "  → 1–2 nm (5–10사이클) 정밀 제어\n"
     "• 비정질(amorphous) VOₓ 선호\n"
     "• V₂O₅(내층)/TiO₂(외층) 이중층:\n"
     "  TiO₂ 외층이 전해질 직접 접촉 차단\n"
     "  Rct: 688 Ω → 202 Ω (1/3 수준)"),
    ("해결 ③\n3단계 스케일업 + 공압 이송",
     "• 공압 이송: 기류로 분말 상시 분산\n"
     "  → 뭉침 방지\n"
     "• 구간 분리형 반응기: 교차 오염 방지\n"
     "• 3단계 스케일업:\n"
     "  FBR(≤10%) → OPERIO 10(≤7%)\n"
     "  → 양산(≤5% @ 10 kg/h)"),
]

for i in range(3):
    x = col_xs[i] + 0.1
    bw = col_w - 0.2

    # ── 난제 박스 (빨강)
    add_rect(s3, x, 1.9, bw, 0.65, C_RED)
    add_text_box(s3, x+0.05, 1.93, bw-0.1, 0.62,
                 problem_data[i][0], 10, True, C_WHITE, PP_ALIGN.CENTER)

    add_rect(s3, x, 2.58, bw, 2.9, C_LIGHTRED, C_RED, 0.5)
    add_text_box(s3, x+0.1, 2.65, bw-0.2, 2.8,
                 problem_data[i][1], 9.5, False, C_DARKGRAY, PP_ALIGN.LEFT)

    # 화살표 (빨강→녹색)
    add_arrow_down(s3, x + bw/2, 5.55, 6.25, C_NAVY, 0.55)

    # ── 해결 박스 (녹색)
    add_rect(s3, x, 6.3, bw, 0.65, C_GREEN)
    add_text_box(s3, x+0.05, 6.33, bw-0.1, 0.62,
                 solution_data[i][0], 10, True, C_WHITE, PP_ALIGN.CENTER)

    add_rect(s3, x, 7.0, bw, 3.5, C_LIGHTGREEN, C_GREEN, 0.5)
    add_text_box(s3, x+0.1, 7.1, bw-0.2, 3.35,
                 solution_data[i][1], 9.5, False, C_DARKGRAY, PP_ALIGN.LEFT)

# 담당기관 뱃지
orgs = ["기계연 + 대성기계", "기계연 (KIMM)", "대성기계 (OPERIO 10)"]
kpis = ["GPC 재현성 ≤5%", "계면저항 ≤50 Ω·cm²", "두께 편차 ≤5%"]
for i in range(3):
    x = col_xs[i] + 0.1
    bw = col_w - 0.2
    add_rect(s3, x, 10.55, bw/2-0.05, 0.5, C_BLUE)
    add_text_box(s3, x+0.02, 10.57, bw/2-0.1, 0.46,
                 orgs[i], 8, True, C_WHITE, PP_ALIGN.CENTER)
    add_rect(s3, x+bw/2+0.05, 10.55, bw/2-0.05, 0.5, C_ORANGE)
    add_text_box(s3, x+bw/2+0.07, 10.57, bw/2-0.1, 0.46,
                 kpis[i], 8, True, C_WHITE, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────
# SLIDE 4 — 그래프 ① Rct 비교 Bar Chart + 이중층 구조 모식도
# ─────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
SW, SH = setup_slide_background(s4, prs)
add_slide_title(s4, prs,
                "난제 ② 해결: V₂O₅/TiO₂ 이중층 ALD 설계",
                "Rct 계면저항 1/3 감소 — 전기화학적 활성 차단 + 이종접합 시너지")

# ── 왼쪽: Rct 비교 Bar Chart ────────────────────────────────────
fig_rct, ax = plt.subplots(figsize=(6, 4.5))
labels = ['Al₂O₃\n단독', 'TiO₂\n단독', 'V₂O₅\n단독', 'V₂O₅/TiO₂\n이중층\n(본 과제)']
values = [323, 380, 688, 202]
colors_bar = ['#4472C4', '#4472C4', '#C00000', '#375623']
bars = ax.bar(labels, values, color=colors_bar, width=0.5, edgecolor='white', linewidth=0.8)

# 값 레이블
for bar, val in zip(bars, values):
    ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 12,
            f'{val} Ω', ha='center', va='bottom', fontsize=11, fontweight='bold',
            color='#C00000' if val == 688 else '#375623' if val == 202 else '#333333')

# 강조 화살표 + 텍스트
ax.annotate('', xy=(3, 220), xytext=(2, 700),
            arrowprops=dict(arrowstyle='->', color='#C00000', lw=2))
ax.text(2.6, 720, '1/3\n수준!', ha='center', fontsize=13, color='#C00000', fontweight='bold')

ax.set_ylabel('계면저항 Rct (Ω)', fontsize=11, fontweight='bold')
ax.set_title('코팅 소재별 계면저항(Rct) 비교\n[Kim et al., ACS AMI, 2024]',
             fontsize=11, fontweight='bold', pad=8)
ax.set_ylim(0, 820)
ax.yaxis.grid(True, alpha=0.3, linestyle='--')
ax.set_facecolor('#FAFAFA')
fig_rct.patch.set_facecolor('white')
fig_rct.tight_layout()

add_fig(s4, fig_rct, 0.5, 1.8, 14.5, 10.5)

# ── 오른쪽: 이중층 구조 모식도 ─────────────────────────────────────
fig_struct, ax2 = plt.subplots(figsize=(5, 5.5))
ax2.set_xlim(0, 10)
ax2.set_ylim(0, 11)
ax2.axis('off')

# 레이어 블록들
layers = [
    (0.5, 9.0, 9.0, 1.6, '#E8E8FF', '#4040A0', '전해질\n(HF, 분해 생성물)',           '#2020A0', 11),
    (0.5, 7.2, 9.0, 1.6, '#DDEEDD', '#205020', 'TiO₂  ~1–2 nm  (외층)\n'
     '• 전해질 직접 접촉 차단\n• Rct 최소화 (이종접합)\n• V₂O₅ 활성 외부 차단', '#205020', 9.5),
    (0.5, 5.2, 9.0, 1.8, '#FFEECC', '#804000', 'V₂O₅  ~1 nm  (내층)\n'
     '• Na 잔류물 제거\n• 산소 방출 억제\n• 혼합 이온·전자 전도(MIEC)', '#804000', 9.5),
    (0.5, 3.2, 9.0, 1.8, '#CCDDFF', '#003090', 'High-Ni SIB 양극재\nNa[Ni₀.₆Co₀.₂Mn₀.₂]O₂',
     '#003090', 11),
]
for (lx, ly, lw, lh, fc, ec, txt, tc, fs) in layers:
    rect = plt.Rectangle((lx, ly), lw, lh, facecolor=fc, edgecolor=ec, linewidth=2)
    ax2.add_patch(rect)
    ax2.text(lx + lw/2, ly + lh/2, txt, ha='center', va='center',
             fontsize=fs, color=tc, fontweight='bold' if fs >= 11 else 'normal',
             multialignment='center')

# 이종접합 계면 마커
ax2.annotate('', xy=(10.2, 7.2), xytext=(10.2, 9.0),
             arrowprops=dict(arrowstyle='<->', color='#C00000', lw=1.5))
ax2.text(10.4, 8.1, '이종접합\nRct 1/3↓', fontsize=8.5, color='#C00000',
         va='center', fontweight='bold')

# 보호 기능 범례
ax2.annotate('', xy=(10.2, 5.2), xytext=(10.2, 7.0),
             arrowprops=dict(arrowstyle='<->', color='#805000', lw=1.5))
ax2.text(10.4, 6.1, 'Na잔류물\n제거', fontsize=8.5, color='#805000',
         va='center', fontweight='bold')

ax2.set_title('V₂O₅(내층)/TiO₂(외층) ALD 이중층\n구조 및 기능 모식도',
              fontsize=11, fontweight='bold', pad=6)
fig_struct.tight_layout()

add_fig(s4, fig_struct, 15.5, 1.8, 17.5, 12.0)

# 핵심 메시지 박스
add_rect(s4, 0.4, SH-1.3, SW-0.8, 1.0, C_NAVY)
add_text_box(s4, 0.6, SH-1.25, SW-1.2, 0.85,
    "★ SIB 양극재에 V₂O₅/TiO₂ ALD 이중층 적용 — 문헌 선례 없음 (세계 최초 시도)  "
    "|  ALD Å 수준 제어로 1–2 nm 유지 → 전기화학적 활성 위험 완전 차단",
    10, True, C_YELLOW, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────
# SLIDE 5 — 그래프 ② RTD 알고리즘 + 3단계 스케일업 로드맵
# ─────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
SW, SH = setup_slide_background(s5, prs)
add_slide_title(s5, prs,
                "난제 ①③ 해결: RTD 알고리즘 & 3단계 스케일업",
                "배치→연속 레시피 이전 방법론 (세계 최초) + 두께 편차 단계별 감소")

# ── 왼쪽: RTD 개념 그림 ─────────────────────────────────────────
fig_rtd, axr = plt.subplots(figsize=(6.5, 5))
t = np.linspace(0, 5, 300)
# RTD 곡선 (감마분포 유사)
mu, sigma = 2.0, 0.6
rtd = (1/(sigma * np.sqrt(2*np.pi))) * np.exp(-0.5*((t-mu)/sigma)**2)
rtd = rtd / rtd.max()

axr.plot(t, rtd, '-', color='#2E75B6', linewidth=2.5, label='체류시간 분포(RTD)')
axr.fill_between(t, rtd, alpha=0.15, color='#2E75B6')

# μ - 2σ 선
t_min = mu - 2*sigma
axr.axvline(t_min, color='#C00000', linewidth=2, linestyle='--',
            label=f'최소 체류시간 (μ-2σ = {t_min:.1f}s)')
# μ 선
axr.axvline(mu, color='#375623', linewidth=1.5, linestyle=':',
            label=f'평균 체류시간 (μ = {mu:.1f}s)')

# t_pulse 마커
t_pulse = 0.6
axr.axvline(t_pulse, color='#ED7D31', linewidth=2, linestyle='-.',
            label=f'배치식 Pulse 시간 (t_pulse = {t_pulse}s)')
axr.fill_betweenx([0, 1.05], 0, t_pulse, alpha=0.12, color='#ED7D31')

# 조건 텍스트
axr.text(0.5, 0.88, 'μ - 2σ ≥ t_pulse\n→ 조건 만족 ✓', fontsize=11,
         color='#C00000', fontweight='bold', ha='center',
         bbox=dict(boxstyle='round,pad=0.3', facecolor='#FFF0F0', edgecolor='#C00000'))

axr.set_xlabel('체류 시간 (s)', fontsize=11)
axr.set_ylabel('RTD (정규화)', fontsize=11)
axr.set_title('RTD 기반 배치→연속식 레시피 이전 알고리즘\n'
              '(분말 최소 체류시간 ≥ 배치식 Pulse 시간 조건)',
              fontsize=10.5, fontweight='bold', pad=8)
axr.legend(fontsize=8.5, loc='upper right')
axr.set_ylim(0, 1.15)
axr.set_xlim(0, 5)
axr.yaxis.grid(True, alpha=0.3, linestyle='--')
axr.set_facecolor('#FAFAFA')
fig_rtd.tight_layout()

add_fig(s5, fig_rtd, 0.5, 2.0, 16.0, 11.5)

# ── 오른쪽: 3단계 스케일업 ──────────────────────────────────────
fig_scale, axs = plt.subplots(figsize=(6, 5))
stages = ['Stage 1\n1차년도\nFBR', 'Stage 2\n2차년도\nOPERIO 10\n(Pilot)', 'Stage 3\n3차년도\n양산형\n연속식']
throughputs = [0.1, 5, 10]          # kg/h 스케일 (시각적)
deviations  = [10, 7, 5]            # 두께 편차 %
bar_colors = ['#4472C4', '#2E75B6', '#1F4E79']

x_pos = np.arange(len(stages))
bars_s = axs.bar(x_pos, throughputs, color=bar_colors, width=0.5,
                 edgecolor='white', linewidth=1)

# 처리량 레이블
throughput_labels = ['수g', 'kg급', '10 kg/h\n★']
for bar, lbl, dev, tp in zip(bars_s, throughput_labels, deviations, throughputs):
    axs.text(bar.get_x() + bar.get_width()/2,
             bar.get_height() + 0.2,
             f'처리량\n{lbl}', ha='center', va='bottom',
             fontsize=9.5, fontweight='bold', color='#1F4E79')
    axs.text(bar.get_x() + bar.get_width()/2,
             bar.get_height() / 2,
             f'두께편차\n≤{dev}%', ha='center', va='center',
             fontsize=10, fontweight='bold', color='white')

axs.set_xticks(x_pos)
axs.set_xticklabels(stages, fontsize=9.5)
axs.set_ylabel('처리량 (kg/h, 개략)', fontsize=10)
axs.set_title('3단계 스케일업 로드맵\n(두께 편차 목표값 단계적 강화)',
              fontsize=10.5, fontweight='bold', pad=8)
axs.set_ylim(0, 13)
axs.yaxis.grid(True, alpha=0.3, linestyle='--')
axs.set_facecolor('#FAFAFA')

# 마일스톤 화살표
for i in range(1, 3):
    axs.annotate('', xy=(i, throughputs[i]*0.8),
                 xytext=(i-1, throughputs[i-1]*0.8+0.5),
                 arrowprops=dict(arrowstyle='->', color='#ED7D31', lw=1.5))

fig_scale.tight_layout()
add_fig(s5, fig_scale, 17.5, 2.0, 15.5, 11.5)

# 하단 수식 박스
add_rect(s5, 0.4, SH-2.3, SW-0.8, 1.9, C_LIGHTBLUE, C_BLUE, 1)
add_text_box(s5, 0.7, SH-2.2, SW-1.4, 0.9,
    "RTD 변환 알고리즘",
    11, True, C_NAVY, PP_ALIGN.LEFT)
add_text_box(s5, 0.7, SH-1.5, SW-1.4, 1.3,
    "배치식 변수  →  연속식 적용:  "
    "반응온도(°C) · Precursor Pulse (s) · Purge (s) · GPC (Å/cyc)  "
    "⟹  최소 체류시간 μ-2σ, 구간 온도, 퍼지 길이/유량, 두께 예측값",
    9.5, False, C_NAVY, PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────
# SLIDE 6 — 그래프 ③ 코팅 소재별 성능 비교 Bar Chart
# ─────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(blank_layout)
SW, SH = setup_slide_background(s6, prs)
add_slide_title(s6, prs,
                "문헌 기반 코팅 소재별 성능 비교",
                "V₂O₅/TiO₂ 이중층의 독창성 — 문헌 선례 없음 (연구 공백 = 특허 선점 기회)")

# ── 왼쪽: 용량 유지율 Bar Chart ─────────────────────────────────
fig_cyc, ax_cyc = plt.subplots(figsize=(6.5, 4.8))
coat_labels = ['Al₂O₃\nALD\n(SIB)', 'TiO₂\nALD\n(SIB)', 'TiO₂/Al₂O₃\n이중층\n(LIB)',
               'V₂O₅\n습식\n(LIB)', 'V₂O₅/TiO₂\n이중층\n(본과제)']
retention = [75, 78, 90.4, 83, None]
cap_vals  = [151, 145, 146, 163, None]
bar_clrs  = ['#4472C4', '#70AD47', '#2E75B6', '#ED7D31', '#C0C0C0']

xs = np.arange(len(coat_labels))
bars_ret = ax_cyc.bar(xs, [v if v else 0 for v in retention],
                      color=bar_clrs, width=0.55, edgecolor='white')

for bar, val in zip(bars_ret, retention):
    if val is not None:
        ax_cyc.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.8,
                    f'{val}%', ha='center', va='bottom', fontsize=10.5, fontweight='bold')

# 마지막 바 (연구공백) 해치
ax_cyc.bar(xs[-1], 85, color='none', hatch='////', edgecolor='#808080',
           width=0.55, label='본 과제\n수행 예정')
ax_cyc.text(xs[-1], 45, '연구\n공백\n→ 본 과제\n수행 예정',
            ha='center', va='center', fontsize=9, color='#808080', fontweight='bold')

ax_cyc.set_xticks(xs)
ax_cyc.set_xticklabels(coat_labels, fontsize=9)
ax_cyc.set_ylabel('용량 유지율 (%)', fontsize=11)
ax_cyc.set_title('코팅 소재별 용량 유지율 비교\n(각 문헌 최적 사이클 기준)',
                 fontsize=10.5, fontweight='bold')
ax_cyc.set_ylim(0, 115)
ax_cyc.axhline(80, color='red', linestyle='--', linewidth=1.2, alpha=0.6,
               label='목표 80% (200cyc)')
ax_cyc.legend(fontsize=8.5)
ax_cyc.yaxis.grid(True, alpha=0.3)
ax_cyc.set_facecolor('#FAFAFA')
fig_cyc.tight_layout()
add_fig(s6, fig_cyc, 0.5, 2.0, 16.0, 11.5)

# ── 오른쪽: 레이더 차트 (소재별 기능 비교) ─────────────────────
fig_radar, ax_r = plt.subplots(figsize=(5.5, 5), subplot_kw=dict(polar=True))
categories = ['HF\n차단', 'Na잔류물\n제거', '이온\n전도성', '전자\n전도성', '상전이\n억제', '산소방출\n억제']
N = len(categories)
angles = [n / float(N) * 2 * np.pi for n in range(N)]
angles += angles[:1]

data_mat = {
    'Al₂O₃': [5, 1, 1, 1, 3, 3],
    'TiO₂':  [3, 1, 3, 2, 3, 3],
    'Al₂O₃/TiO₂': [5, 1, 3, 2, 4, 3],
    'V₂O₅(예측)': [3, 5, 5, 4, 3, 5],
}
colors_r = ['#4472C4', '#70AD47', '#2E75B6', '#C00000']
for (name, vals), clr in zip(data_mat.items(), colors_r):
    vals_plot = vals + vals[:1]
    lw = 2.5 if name == 'V₂O₅(예측)' else 1.5
    ls = '-' if name == 'V₂O₅(예측)' else '--'
    ax_r.plot(angles, vals_plot, 'o-', color=clr, linewidth=lw, linestyle=ls, label=name)
    if name == 'V₂O₅(예측)':
        ax_r.fill(angles, vals_plot, alpha=0.12, color=clr)

ax_r.set_xticks(angles[:-1])
ax_r.set_xticklabels(categories, fontsize=9.5)
ax_r.set_ylim(0, 5.5)
ax_r.set_yticks([1,2,3,4,5])
ax_r.set_yticklabels(['','','','',''], fontsize=7)
ax_r.set_title('코팅 소재별 기능 비교\n(레이더 차트)',
               fontsize=10.5, fontweight='bold', pad=15)
ax_r.legend(loc='upper right', bbox_to_anchor=(1.35, 1.1), fontsize=8.5)
ax_r.set_facecolor('#FAFAFA')
fig_radar.tight_layout()
add_fig(s6, fig_radar, 17.5, 2.0, 15.5, 12.0)

# 메시지
add_rect(s6, 0.4, SH-1.4, SW-0.8, 1.1, C_NAVY)
add_text_box(s6, 0.6, SH-1.32, SW-1.2, 0.95,
    "★ V₂O₅/TiO₂ 이중층 ALD의 SIB 양극재 적용 — 선행 연구 없음  "
    "|  V₂O₅의 'Na 잔류물 제거 + MIEC' 기능을 ALD 두께 제어로 안전하게 활용하는 최초 시도",
    10, True, C_YELLOW, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────
# SLIDE 7 — 요약 테이블 + 핵심 메시지
# ─────────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(blank_layout)
SW, SH = setup_slide_background(s7, prs)
add_slide_title(s7, prs,
                "3대 난제 × 해결방안 요약 및 본 컨소시엄의 차별성",
                "기계연 + 대성기계만이 이 3가지 난제를 동시에 해결할 수 있는 이유")

# ── 요약 테이블 (직접 그리기) ──────────────────────────────────
table_data = [
    ["", "난제", "핵심 원인", "기술적 해결책", "담당 기관", "검증 지표(3차)"],
    ["①", "배치→연속 이전 불가", "공정 환경 근본 차이\n(정지↔비행반응)",
     "RTD 정량 변환 알고리즘\nμ-2σ ≥ t_pulse",
     "기계연\n+ 대성기계", "GPC 재현성\n≤5%"],
    ["②", "V₂O₅ 전기화학 활성", "전압 범위 중복\n(2.0–3.5V / 2.0–4.2V)",
     "ALD 두께 1–2 nm 제어\n+ 비정질 VOₓ + TiO₂ 외층",
     "기계연\n(KIMM)", "계면저항\n≤50 Ω·cm²"],
    ["③", "양산 균일도 미확보", "분말 뭉침·정체구간\n(처리량 ×3,000배)",
     "공압이송 분산\n+ 구간분리 반응기",
     "대성기계\n(OPERIO 10)", "두께 편차\n≤5% @10kg/h"],
]

# 헤더 색상
hdr_fill  = C_NAVY
body_fills = [C_LIGHTRED, RGBColor(0xFF,0xF0,0xCC), C_LIGHTGREEN]
hdr_widths = [1.2, 5.5, 6.0, 8.0, 4.5, 5.0]
total_w = sum(hdr_widths)  # 30.2 cm
scale = (SW - 0.8) / total_w
hdr_w_scaled = [w * scale for w in hdr_widths]

tbl_x = 0.4
tbl_y = 1.85
row_h = [0.65, 1.6, 1.8, 1.6]

for ri, row in enumerate(table_data):
    cx = tbl_x
    for ci, cell in enumerate(row):
        cw = hdr_w_scaled[ci]
        ch = row_h[ri]
        if ri == 0:
            fill = hdr_fill
            tc = C_WHITE
            fs = 10; bold = True
        else:
            fill = body_fills[ri-1] if ci > 0 else RGBColor(0xEE, 0xEE, 0xFF)
            tc = C_DARKGRAY if ci > 0 else C_NAVY
            fs = 9.5; bold = (ci == 0)

        add_rect(s7, cx, tbl_y, cw, ch, fill,
                 RGBColor(0xAA, 0xAA, 0xAA), 0.5)
        add_text_box(s7, cx+0.05, tbl_y+0.05, cw-0.1, ch-0.1,
                     cell, fs, bold, tc, PP_ALIGN.CENTER)
        cx += cw
    tbl_y += row_h[ri]

# ── 3개 핵심 강점 박스 ─────────────────────────────────────────
bx = 0.4
by = tbl_y + 0.35
bw = (SW - 0.8) / 3 - 0.2
strengths = [
    ("기계연 (KIMM)\n20년+ ALD 연구 역량",
     "• Science 2009 (피인용 503회)\n"
     "• ZnF₂ ALD 배터리 적용 (2025)\n"
     "• Å 수준 두께 제어 기술 내재화\n"
     "• 배치식 ALD 장비 2대 운용"),
    ("대성기계 OPERIO 10\n국내 유일 연속식 분말 ALD",
     "• 공압 이송 기반 상압 연속 공정\n"
     "• RTD 정밀 제어 (국내 선도)\n"
     "• Turn-key ALD 장비 공급 경험\n"
     "• 10 kg/h 양산 장비 보유"),
    ("협력 모델 (세계 최초)\n배치↔연속 RTD 이전 방법론",
     "• 기계연 레시피 → 대성기계 이전\n"
     "• 3단계 스케일업 리스크 분산\n"
     "• 코팅 시료 ↔ 분석 피드백 루프\n"
     "• SIB 양극재 V₂O₅/TiO₂ 독창성"),
]
fills_str = [C_LIGHTBLUE, C_LIGHTGREEN, RGBColor(0xFF, 0xF4, 0xCC)]
borders_str = [C_BLUE, C_GREEN, C_ORANGE]
for i, (title, body) in enumerate(strengths):
    x_s = bx + i * (bw + 0.3)
    add_rect(s7, x_s, by, bw, 0.65, borders_str[i])
    add_text_box(s7, x_s+0.05, by+0.05, bw-0.1, 0.58,
                 title, 10, True, C_WHITE, PP_ALIGN.CENTER)
    add_rect(s7, x_s, by+0.68, bw, 2.8, fills_str[i],
             borders_str[i], 0.75)
    add_text_box(s7, x_s+0.1, by+0.78, bw-0.2, 2.6,
                 body, 9.5, False, C_DARKGRAY, PP_ALIGN.LEFT)

# 결론 배너
add_rect(s7, 0.4, SH-1.0, SW-0.8, 0.75, C_NAVY)
add_text_box(s7, 0.6, SH-0.95, SW-1.2, 0.65,
    "★ 본 컨소시엄(기계연 + 대성기계)만이 ALD 레시피 정밀 개발 + 연속식 10 kg/h 양산 이전을 동시에 구현할 수 있는 국내 유일 협력 모델",
    10.5, True, C_YELLOW, PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────
# SLIDE 8 — 검증 지표 달성 경로 (KPI 정리)
# ─────────────────────────────────────────────────────────────────
s8 = prs.slides.add_slide(blank_layout)
SW, SH = setup_slide_background(s8, prs)
add_slide_title(s8, prs,
                "ALD 파트 연차별 KPI 및 달성 경로",
                "3가지 핵심 검증 지표 — 단계별 엄격화 구조")

# KPI 꺾은선 그래프
fig_kpi, axes = plt.subplots(1, 3, figsize=(14, 5))

kpi_data = [
    ("두께 편차 (%)\n[목표: 낮을수록 좋음]",
     [10, 7, 5], '≤', '%', '#C00000', True),
    ("초기 쿨롱 효율 (%)\n[목표: 높을수록 좋음]",
     [80, 83, 85], '≥', '%', '#375623', False),
    ("계면저항 Rct (Ω·cm²)\n[목표: 낮을수록 좋음]",
     [70, 60, 50], '≤', ' Ω·cm²', '#C00000', True),
]
years = ['1차년도', '2차년도', '3차년도']

for ax_k, (title, vals, sym, unit, clr, lower_better) in zip(axes, kpi_data):
    ax_k.plot(years, vals, 'o-', color=clr, linewidth=2.5,
              markersize=9, markerfacecolor=clr, zorder=3)
    for xi, (yr, v) in enumerate(zip(years, vals)):
        offset = -3 if lower_better else 1.5
        ax_k.annotate(f'{sym}{v}{unit}',
                      xy=(xi, v), xytext=(0, 12 if not lower_better else -18),
                      textcoords='offset points',
                      ha='center', fontsize=11, fontweight='bold', color=clr)

    # 목표선 (최종값)
    ax_k.axhline(vals[-1], color=clr, linestyle=':', linewidth=1.2, alpha=0.5)

    # 배경 그라데이션 느낌
    ax_k.fill_between(years, vals, [vals[0]]*3 if lower_better else [0]*3,
                      alpha=0.08, color=clr)
    ax_k.set_title(title, fontsize=10.5, fontweight='bold', pad=8)
    ax_k.set_xticks(range(3))
    ax_k.set_xticklabels(years, fontsize=10)
    ax_k.yaxis.grid(True, alpha=0.3, linestyle='--')
    ax_k.set_facecolor('#FAFAFA')

    # y축 범위 조정
    mn, mx = min(vals), max(vals)
    margin = (mx - mn) * 0.4
    if lower_better:
        ax_k.set_ylim(mn - margin, mx + margin*1.5)
    else:
        ax_k.set_ylim(mn - margin*1.5, mx + margin)

    # 화살표로 방향 표시
    arrow_y = vals[-1]
    direction = '↓ 감소' if lower_better else '↑ 향상'
    ax_k.text(0.02, 0.05 if lower_better else 0.95,
              direction, transform=ax_k.transAxes,
              fontsize=12, color=clr, fontweight='bold',
              va='bottom' if lower_better else 'top')

fig_kpi.suptitle('연차별 핵심 KPI 달성 로드맵', fontsize=13, fontweight='bold', y=1.02)
fig_kpi.tight_layout(w_pad=3)
add_fig(s8, fig_kpi, 0.5, 2.0, SW-1, 12.5)

# 측정 기준 주석
add_rect(s8, 0.4, SH-1.5, SW-0.8, 1.2, C_LIGHTBLUE, C_BLUE, 0.75)
add_text_box(s8, 0.6, SH-1.45, SW-1.2, 1.1,
    "측정 기준:\n"
    "• 두께 편차: 1차=배치식 FBR / 2·3차=10 kg/h 연속 운전 기준  "
    "• 쿨롱 효율: 1차=coin cell / 3차=1Ah pouch cell  "
    "• 계면저항: EIS 평가, 200사이클 후 25°C",
    9, False, C_NAVY, PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────
# 저장
# ─────────────────────────────────────────────────────────────────
output_path = r"d:\23_GitHub\01_Office_PC\030_Project\032_Proposed Project\08_2026년 첨단소재 사업화 기술 개발(성과확장형)\ALD_기술난제_해결방안.pptx"
prs.save(output_path)
print(f"✅ PPT 저장 완료: {output_path}")
print(f"   총 {len(prs.slides)} 슬라이드")
for i, slide in enumerate(prs.slides, 1):
    print(f"   Slide {i}: {len(slide.shapes)} shapes")
