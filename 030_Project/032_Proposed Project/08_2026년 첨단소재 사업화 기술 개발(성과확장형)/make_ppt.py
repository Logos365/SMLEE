"""
2026 첨단소재 사업화 기술 개발(성과확장형) — 2세부 선정평가 발표자료
페이지별 콘텐츠 가이드 기반 PPT 생성 스크립트
대상 페이지: 3, 5, 6, 8, 9, 12, 15, 24
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── 컬러 팔레트 ──────────────────────────────────────────────────
NAVY    = RGBColor(0x1B, 0x3A, 0x6B)   # 헤더 배너
BLUE    = RGBColor(0x2E, 0x5F, 0xA3)   # 기계연
BLUE2   = RGBColor(0x2E, 0x86, 0xC1)   # 보조 파랑
AMBER   = RGBColor(0xF0, 0xA5, 0x00)   # 마일스톤 / 강조
RED     = RGBColor(0xC8, 0x10, 0x2E)   # 독창 강조
GREEN   = RGBColor(0x27, 0xAE, 0x60)   # 대성기계 / 달성
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xF4, 0xF6, 0xFA)   # 슬라이드 배경
LGRAY   = RGBColor(0xE8, 0xEC, 0xF2)   # 셀 배경
DGRAY   = RGBColor(0x44, 0x44, 0x55)   # 본문 텍스트


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header(slide, title, subtitle=None):
    """상단 네이비 배너 + 제목"""
    sh = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Cm(0), Cm(0), Cm(33.87), Cm(2.4)
    )
    sh.fill.solid(); sh.fill.fore_color.rgb = NAVY
    sh.line.fill.background()

    tf = sh.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(20)
    run.font.color.rgb = WHITE
    run.font.name = "맑은 고딕"
    tf.margin_left = Cm(0.5)
    tf.margin_top = Cm(0.3)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(11)
        r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
        r2.font.name = "맑은 고딕"


def add_box(slide, left, top, width, height, text, font_size=11,
            bg_color=WHITE, text_color=DGRAY, bold=False, border_color=NAVY,
            align=PP_ALIGN.LEFT, border_width=Pt(1)):
    sh = slide.shapes.add_shape(1, left, top, width, height)
    sh.fill.solid(); sh.fill.fore_color.rgb = bg_color
    sh.line.color.rgb = border_color
    sh.line.width = border_width
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left  = Cm(0.25)
    tf.margin_right = Cm(0.15)
    tf.margin_top   = Cm(0.1)
    tf.margin_bottom= Cm(0.1)
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color
        run.font.bold = bold
        run.font.name = "맑은 고딕"
    return sh


def add_label(slide, left, top, text, font_size=10, color=DGRAY, bold=False):
    tb = slide.shapes.add_textbox(left, top, Cm(30), Cm(0.8))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "맑은 고딕"


def add_table(slide, left, top, width, height, headers, rows,
              header_bg=NAVY, row_bg=WHITE, alt_bg=LGRAY,
              font_size=9, header_font_size=10):
    n_cols = len(headers)
    n_rows = len(rows)
    tbl = slide.shapes.add_table(n_rows + 1, n_cols, left, top, width, height).table

    col_w = width // n_cols
    for i in range(n_cols):
        tbl.columns[i].width = col_w

    row_h = height // (n_rows + 1)
    for i in range(n_rows + 1):
        tbl.rows[i].height = row_h

    # 헤더
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = hdr
        run.font.bold = True
        run.font.size = Pt(header_font_size)
        run.font.color.rgb = WHITE
        run.font.name = "맑은 고딕"

    # 데이터
    for ri, row in enumerate(rows):
        bg = alt_bg if ri % 2 == 1 else row_bg
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(font_size)
            run.font.color.rgb = DGRAY
            run.font.name = "맑은 고딕"
            if str(val).startswith("★") or "★" in str(val):
                run.font.bold = True
                run.font.color.rgb = RED


def add_page_num(slide, num):
    tb = slide.shapes.add_textbox(Cm(32), Cm(18.5), Cm(1.5), Cm(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0xAA, 0xAA, 0xBB)
    run.font.name = "맑은 고딕"


# ════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Cm(33.87)   # 16:9
prs.slide_height = Cm(19.05)
blank = prs.slide_layouts[6]   # 완전 빈 슬라이드


# ══════════════════════════════════════════════════════════════════
# PAGE 3 — ALD 코팅 파트
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
add_header(sl, "ALD(원자층 증착) 기반 양극재 계면 안정화 기술",
           "SIB 고-Ni 양극재 코팅 소재 선정 근거 및 독창 구조")
add_page_num(sl, 3)

# ── 왼쪽: ALD 공정 원리
add_box(sl, Cm(0.4), Cm(2.6), Cm(9.5), Cm(0.6),
        "■ ALD 공정 원리", font_size=11, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)
add_box(sl, Cm(0.4), Cm(3.2), Cm(9.5), Cm(3.2),
        "전구체 A Pulse → Purge → 반응제(H₂O) Pulse → Purge  반복\n"
        "자기제한 반응(Self-limiting): GPC = 0.1~2.0 Å/cycle\n"
        "→ Å 수준 두께 제어 + 완전 균일 피복(conformal coating)",
        font_size=10, bg_color=WHITE, border_color=BLUE2)

# ── 왼쪽 하단: 코팅 구조 모식도
add_box(sl, Cm(0.4), Cm(6.5), Cm(9.5), Cm(0.6),
        "■ V₂O₅/TiO₂ 이중층 구조 모식도", font_size=11, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)
layers = [
    ("전해질 (HF 접촉 차단됨)",           BG,    DGRAY, False),
    ("TiO₂  ~1–2 nm  (외층)",           RGBColor(0xCF,0xE2,0xFF), BLUE, True),
    ("✦ 이종접합 계면 → Rct 1/3 감소 (202 Ω)", RGBColor(0xFF,0xEC,0xEC), RED,  True),
    ("V₂O₅  ~1 nm  (내층)",             RGBColor(0xFF,0xF0,0xCC), RGBColor(0x8B,0x6A,0x00), True),
    ("SIB 고-Ni 양극재  NaNi₀.₈Mn₀.₁Co₀.₁O₂", NAVY,  WHITE, True),
]
y = 7.1
for txt, bg, fc, bld in layers:
    h = 1.3 if "NaNi" in txt else (1.0 if "이종접합" in txt else 0.85)
    add_box(sl, Cm(0.4), Cm(y), Cm(9.5), Cm(h),
            txt, font_size=10, bg_color=bg, text_color=fc, bold=bld,
            border_color=BLUE2, align=PP_ALIGN.CENTER)
    y += h

# ── 중앙: 코팅 소재별 역할 비교표
add_box(sl, Cm(10.2), Cm(2.6), Cm(14.5), Cm(0.6),
        "■ 코팅 소재별 역할 비교", font_size=11, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)
add_table(sl,
    Cm(10.2), Cm(3.2), Cm(14.5), Cm(6.0),
    ["소재", "주 기능", "두께", "본 과제 적용"],
    [
        ["Al₂O₃",                "HF 차단, AlF₃ 보호층 형성",         "1–3 nm",   "1차년도"],
        ["TiO₂",                 "상전이 억제, 이온 전도성",           "1–3 nm",   "1차년도"],
        ["★Al₂O₃/TiO₂ 이중층",   "HF 차단 + 상전이 억제 복합",        "3–5 nm",  "1·2차년도 ★"],
        ["V₂O₅",                 "Na 잔류물 제거, 이온·전자 전도성",   "1–2 nm",   "2차년도"],
        ["★V₂O₅/TiO₂ 이중층",    "Na 잔류물 제거 + Rct 최소화",        "2–3 nm",  "3차년도 ★★"],
    ],
    font_size=9, header_font_size=10
)

# ── 오른쪽: 문헌 근거 강조박스
add_box(sl, Cm(25.0), Cm(2.6), Cm(8.5), Cm(0.6),
        "■ 문헌 기반 성능 근거", font_size=11, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)
add_box(sl, Cm(25.0), Cm(3.2), Cm(8.5), Cm(2.8),
        "Al₂O₃/TiO₂ 이중층\n100사이클 후  90.4%  용량 유지\nChen et al., ACS AMI, 2024",
        font_size=10, bg_color=RGBColor(0xE8,0xF4,0xFF), border_color=BLUE2)
add_box(sl, Cm(25.0), Cm(6.2), Cm(8.5), Cm(3.1),
        "★ V₂O₅/TiO₂ 이중층  (본 과제)\n계면저항  202 Ω\n(단독 V₂O₅ 688 Ω 대비  1/3)\nKim et al., ACS AMI, 2024",
        font_size=10, bg_color=RGBColor(0xFF,0xEE,0xEE),
        text_color=RED, border_color=RED, bold=True)

# ── 하단 강조 배너
add_box(sl, Cm(0.4), Cm(17.0), Cm(33.0), Cm(1.7),
        "SIB 양극재에 V₂O₅/TiO₂ ALD 이중층을 적용한 문헌 선례 없음 → 세계 최초 시도\n"
        "Al₂O₃·TiO₂가 해결 못하는 Na 잔류물 제거 + Rct 최소화를 동시에 달성하는 독창 구조",
        font_size=11, bold=True,
        bg_color=NAVY, text_color=WHITE, border_color=NAVY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# PAGE 5 — 기관 역량
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
add_header(sl, "2세부 핵심 기관 역량 — 기계연 × 대성기계의 분업 체계",
           "배치식 ALD 레시피 개발(기계연) ↔ 연속식 양산 스케일업(대성기계)")
add_page_num(sl, 5)

# 기계연 영역
add_box(sl, Cm(0.4), Cm(2.6), Cm(15.8), Cm(0.7),
        "한국기계연구원 (KIMM)  —  배치식 ALD 레시피 개발 + 코팅층 종합 분석",
        font_size=11, bold=True, bg_color=BLUE, text_color=WHITE, border_color=BLUE)
add_box(sl, Cm(0.4), Cm(3.3), Cm(15.8), Cm(4.5),
        "• 배치식 ALD 장비 2대 동시 운용 — TiO₂, Al₂O₃, V₂O₅, ZnO 병렬 레시피 개발\n"
        "• HR-TEM, EELS, XPS, Raman, FTIR, XRD 등 원자 단위 분석 인프라 완비\n"
        "• ALD 분야 SCI 논문 다수 (Science 2009 포함, 누적 피인용 2,000+ 회)\n"
        "• ZnF₂ ALD 배터리 계면 안정화 논문·특허 보유 (Adv. Sci. 2025)",
        font_size=10, bg_color=WHITE, border_color=BLUE2)
add_table(sl, Cm(0.4), Cm(7.8), Cm(15.8), Cm(3.0),
    ["연차", "한국기계연구원 주요 수행 내용"],
    [
        ["1차년도", "TiO₂·Al₂O₃ ALD 레시피 개발 및 FBR 이전"],
        ["2차년도", "V₂O₅·ZnO ALD 레시피 확장 및 OPERIO 10 이전"],
        ["3차년도", "최종 레시피 DB 완성 + Post-mortem 분석"],
    ],
    header_bg=BLUE, font_size=10
)

# 대성기계 영역
add_box(sl, Cm(17.1), Cm(2.6), Cm(16.3), Cm(0.7),
        "대성기계공업㈜  —  연속식 분말 ALD 운용 + 10 kg/h 양산 실증",
        font_size=11, bold=True, bg_color=GREEN, text_color=WHITE, border_color=GREEN)
add_box(sl, Cm(17.1), Cm(3.3), Cm(16.3), Cm(4.5),
        "• 공압 이송 기반 연속식 분말 ALD 장비 OPERIO 10 보유 (국내 유일)\n"
        "• RTD(체류시간 분포) 기반 공정 정밀 제어 기술\n"
        "• 분말 정량 피딩 및 가스 분산 기술\n"
        "• ALD 장비 Turn-key 상용화 공급 경험",
        font_size=10, bg_color=WHITE, border_color=GREEN)
add_table(sl, Cm(17.1), Cm(7.8), Cm(16.3), Cm(3.0),
    ["연차", "대성기계공업 주요 수행 내용"],
    [
        ["1차년도", "R&D용 FBR 구축 + 기계연 TiO₂/Al₂O₃ 레시피 적용"],
        ["2차년도", "OPERIO 10에 4종 레시피 적용·안정화"],
        ["3차년도", "10 kg/h 연속 운전 실증 (두께 편차 ≤5%)"],
    ],
    header_bg=GREEN, font_size=10
)

# 중앙 연계 구조
add_box(sl, Cm(0.4), Cm(11.0), Cm(33.0), Cm(0.6),
        "■ 기관 연계 구조 (폐루프 피드백 파이프라인)",
        font_size=10, bold=True, bg_color=NAVY, text_color=WHITE, border_color=NAVY)
flow_items = [
    ("기계연 배치식 ALD\n레시피 개발·확립", BLUE),
    ("레시피 패키지 이전 →", BG),
    ("대성기계 연속식 ALD\nRTD 환산 적용·코팅", GREEN),
    ("← 코팅 시료 제공", BG),
    ("기계연 분석 검증\nTEM/XPS/EELS", BLUE),
    ("← 분석 결과 피드백", BG),
    ("레시피 미세 조정·재이전", RGBColor(0xFF,0xF0,0xCC)),
]
x = 0.4
for txt, bg in flow_items:
    is_arrow = "→" in txt or "←" in txt
    fc = AMBER if is_arrow else (WHITE if bg in (BLUE, GREEN) else DGRAY)
    add_box(sl, Cm(x), Cm(11.65), Cm(4.6), Cm(2.5),
            txt, font_size=9, bg_color=bg, text_color=fc,
            bold=not is_arrow, border_color=BLUE2 if not is_arrow else BG,
            align=PP_ALIGN.CENTER)
    x += 4.7

add_box(sl, Cm(0.4), Cm(14.3), Cm(33.0), Cm(1.1),
        "RTD 기반 배치→연속식 변환 방법론 최초 확립  |  국내 유일 배치식↔연속식 분말 ALD 협력 모델",
        font_size=12, bold=True,
        bg_color=RGBColor(0xFF,0xF0,0xCC), text_color=NAVY,
        border_color=AMBER, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# PAGE 6 — 연구추진 체계 및 일정
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
add_header(sl, "3개년 연구추진 체계 및 연차별 일정",
           "기관별 역할 체계도 + 간트차트 + 마일스톤")
add_page_num(sl, 6)

# 역할 체계도
add_box(sl, Cm(0.4), Cm(2.6), Cm(33.0), Cm(0.55),
        "■ 기관별 역할 체계도", font_size=10, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)
orgs = [
    ("전북대학교 (1세부 주관)\n총괄 조정 / 전기화학 평가 / KPI 검증", NAVY, WHITE, Cm(10.5), Cm(1.1)),
    ("한국기계연구원 (2세부)\n배치식 ALD 레시피 개발\n코팅층 종합 분석", BLUE, WHITE, Cm(9.0), Cm(1.5)),
    ("대성기계공업 (2세부)\n연속식 분말 ALD 스케일업\n10 kg/h 양산 실증", GREEN, WHITE, Cm(9.0), Cm(1.5)),
    ("에버인더스\nSIB 양극재 공급", BLUE2, WHITE, Cm(4.5), Cm(1.1)),
    ("씨오알엔\n1Ah 파우치셀 제작", BLUE2, WHITE, Cm(4.5), Cm(1.1)),
]
xs = [Cm(11.2), Cm(1.5), Cm(22.0), Cm(0.4), Cm(18.0)]
ys = [Cm(3.2), Cm(5.0), Cm(5.0), Cm(6.8), Cm(6.8)]
for i, (txt, bg, fc, w, h) in enumerate(orgs):
    add_box(sl, xs[i], ys[i], w, h, txt, font_size=9, bg_color=bg,
            text_color=fc, bold=True, border_color=NAVY, align=PP_ALIGN.CENTER)

# 간트차트
add_box(sl, Cm(0.4), Cm(8.6), Cm(33.0), Cm(0.55),
        "■ 연차별 추진 일정 (간트차트)", font_size=10, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)
add_table(sl, Cm(0.4), Cm(9.15), Cm(33.0), Cm(7.2),
    ["과제 내용", "1차년도", "2차년도", "3차년도"],
    [
        ["[기계연] TiO₂·Al₂O₃ ALD 레시피 개발",         "●●●●●●", "",        ""],
        ["[기계연] V₂O₅·ZnO ALD 레시피 개발",           "",        "●●●●●●", ""],
        ["[기계연] 코팅층 종합 분석 (HRTEM/EELS/XPS)",   "●●●●",   "●●●●●●", "●●●●"],
        ["[기계연] 레시피 이전 검증 및 피드백",           "●●●●",   "●●●●●●", "●●●●"],
        ["[대성기계] FBR 구축 + TiO₂/Al₂O₃ 적용",      "●●●●●●●●", "",      ""],
        ["[대성기계] OPERIO 10 연속 공정 안정화 (4종)",  "",        "●●●●●●●●",""],
        ["[대성기계] 10 kg/h 양산 운전 실증",            "",        "",        "●●●●●●●●"],
        ["전기화학 평가 (Coin→Pouch)",                  "●●",      "●●●●",   "●●●●"],
    ],
    font_size=9, header_font_size=10
)

# 마일스톤
add_box(sl, Cm(0.4), Cm(16.4), Cm(10.8), Cm(1.9),
        "M1  1차년도 말\nTiO₂/Al₂O₃ 레시피 패키지\n→ 대성기계 FBR 이전 완료",
        font_size=9, bg_color=RGBColor(0xFF,0xF0,0xCC),
        text_color=NAVY, border_color=AMBER, bold=True, align=PP_ALIGN.CENTER)
add_box(sl, Cm(11.5), Cm(16.4), Cm(10.5), Cm(1.9),
        "M2  2차년도 말\n4종 레시피 OPERIO 10 이전\n배치식 vs. 연속식 품질 동등성 검증",
        font_size=9, bg_color=RGBColor(0xFF,0xF0,0xCC),
        text_color=NAVY, border_color=AMBER, bold=True, align=PP_ALIGN.CENTER)
add_box(sl, Cm(22.3), Cm(16.4), Cm(11.1), Cm(1.9),
        "M3  3차년도 말\n10 kg/h 24시간 연속 운전 실증\n+ 1Ah 파우치셀 최종 검증",
        font_size=9, bg_color=RGBColor(0xFF,0xF0,0xCC),
        text_color=NAVY, border_color=AMBER, bold=True, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# PAGE 8 — 연속식 ALD 공정 개발
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
add_header(sl, "연속식 분말 ALD 공정 개발 — 배치식 레시피 → 양산 이전 파이프라인",
           "RTD 기반 정량 변환 알고리즘 + 3단계 스케일업 로드맵")
add_page_num(sl, 8)

# 왼쪽: 배치식 vs 연속식 비교표
add_box(sl, Cm(0.4), Cm(2.6), Cm(15.5), Cm(0.6),
        "■ 배치식 vs. 연속식 ALD 비교", font_size=11, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)
add_table(sl, Cm(0.4), Cm(3.2), Cm(15.5), Cm(4.5),
    ["구분", "배치식 FBR (기계연)", "연속식 공압 이송 (대성기계)"],
    [
        ["레시피 개발 용이성", "★★★★★ 매우 높음", "★★☆☆☆ 낮음"],
        ["처리량",            "소량~중간",         "10 kg/h 이상 ★"],
        ["두께 제어",          "정밀 (직접 측정)",  "RTD 기반 간접 제어"],
        ["본 과제 역할",       "레시피 소스",        "양산 적용 플랫폼"],
    ],
    font_size=9
)

# 왼쪽 하단: RTD 핵심 기술
add_box(sl, Cm(0.4), Cm(8.0), Cm(15.5), Cm(0.6),
        "■ 핵심 기술: RTD 기반 레시피 이전 방법론", font_size=11, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)
add_box(sl, Cm(0.4), Cm(8.6), Cm(15.5), Cm(2.5),
        "문제: 배치식(정지 반응) → 연속식(비행 반응) 간 공정 환경 근본적 차이\n"
        "해결: Pulse 시간을 체류시간 분포(RTD) 최솟값으로 변환하는 정량적 알고리즘",
        font_size=10, bg_color=WHITE, border_color=BLUE2)
add_box(sl, Cm(0.4), Cm(11.1), Cm(15.5), Cm(1.4),
        "최소 체류시간:  μ - 2σ  ≥  t_pulse",
        font_size=13, bold=True,
        bg_color=NAVY, text_color=WHITE, border_color=NAVY, align=PP_ALIGN.CENTER)

add_table(sl, Cm(0.4), Cm(12.6), Cm(15.5), Cm(4.2),
    ["배치식 변수 (기계연)", "연속식 변수 (대성기계)"],
    [
        ["반응 온도 (°C)",          "반응 구간 온도 설정값"],
        ["Precursor pulse 시간 (s)","분말 최소 체류시간 (RTD)"],
        ["Purge 시간 (s)",          "퍼지 구간 길이 및 가스 유량"],
        ["GPC (Å/cycle)",          "코팅 두께 예측값"],
        ["공정 윈도우 (온도 범위)", "연속 공정 허용 온도 편차"],
    ],
    font_size=9
)

# 오른쪽: 3단계 스케일업 로드맵
add_box(sl, Cm(16.2), Cm(2.6), Cm(17.2), Cm(0.6),
        "■ 3단계 스케일업 로드맵", font_size=11, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)
stages = [
    ("STAGE 1\n1차년도",
     "배치식 FBR\n기계연 R&D 장비\n소량(g) · ≤10% · TiO₂·Al₂O₃\nSelf-limiting · GPC 포화 곡선",
     BLUE, Cm(3.2)),
    ("STAGE 2\n2차년도",
     "OPERIO 10\n대성기계 파일럿\nkg급 · ≤7% · 4종 이전\nRTD 공정 매핑 · 연속 안정화\n품질 동등성 검증",
     BLUE2, Cm(3.8)),
    ("STAGE 3\n3차년도 ★",
     "양산형 연속식\n대성기계 양산 장비\n10 kg/h ★ · ≤5% ★\n24h 연속 운전 실증\n1Ah 파우치셀 검증\nScale-up 가이드 완성",
     GREEN, Cm(4.5)),
]
y_s = 3.2
for badge, body, bg, h in stages:
    add_box(sl, Cm(16.2), Cm(y_s), Cm(17.2), Cm(0.6),
            badge, font_size=10, bold=True,
            bg_color=bg, text_color=WHITE, border_color=bg, align=PP_ALIGN.CENTER)
    add_box(sl, Cm(16.2), Cm(y_s + 0.6), Cm(17.2), h,
            body, font_size=10, bg_color=WHITE, border_color=bg)
    y_s += h + 0.8

add_box(sl, Cm(16.2), Cm(16.6), Cm(17.2), Cm(1.8),
        "핵심 기술 — RTD 정량 변환\n최소 체류시간: μ - 2σ ≥ t_pulse\n→ 기술 이전 재현성 확보",
        font_size=10, bold=True,
        bg_color=RGBColor(0xFF,0xF0,0xCC), text_color=NAVY,
        border_color=AMBER, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════
# PAGE 9 — 소재 항목 KPI
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
add_header(sl, "연구개발 목표 — 소재 항목 평가지표 (KPI)",
           "연차별 목표값 및 측정 기준")
add_page_num(sl, 9)

add_table(sl, Cm(0.4), Cm(2.8), Cm(33.0), Cm(7.5),
    ["평가항목 (가중치)", "1차년도", "2차년도", "3차년도 (최종)", "측정 기준"],
    [
        ["ALD 코팅 균일도\n두께 편차 (%)  [30%]", "≤10%", "≤7%",  "≤5% ★",  "1차: 배치식 / 2·3차: 10 kg/h 연속"],
        ["초기 쿨롱 효율 (%)  [20%]",               "≥80%", "≥83%", "≥85% ★", "1차: coin cell / 3차: 1Ah pouch cell"],
        ["계면저항 (Ω·cm²)  [20%]",                 "≤70",  "≤60",  "≤50 ★",  "EIS 평가, 사이클 후 25°C"],
        ["코팅층 피복률 (%)  [15%]",                 "≥85%", "≥90%", "≥95% ★", "TEM/XPS 분석 기준"],
        ["GPC 재현성 (%)  [15%]",                   "≤15%", "≤10%", "≤5% ★",  "1차: batch / 3차: 24시간 연속 운전"],
    ],
    font_size=10, header_font_size=11
)

# 점검 의견 박스
notes = [
    ("두께 편차  최종 ≤5%", "배치식·Pilot·양산 단계별로 엄격해지는 구조 —\n단계별 측정 조건 명시 필수"),
    ("초기 쿨롱 효율  최종 ≥85%", "coin cell(1차) → pouch cell(3차) 측정 기준 변화\n명확히 표기 필요"),
    ("계면저항 단위  Ω·cm²", "단순 Ω와 혼동 주의\n슬라이드 표기 통일 확인"),
    ("용량유지율 포함 시", "200사이클 후 ≥80%\n(3차년도 coin cell 기준) 적용"),
]
y_n = 10.8
for i, (title, body) in enumerate(notes):
    x_n = 0.4 + (i % 2) * 16.8
    if i % 2 == 0 and i > 0:
        y_n += 3.6
    add_box(sl, Cm(x_n), Cm(y_n), Cm(16.0), Cm(0.55),
            title, font_size=10, bold=True,
            bg_color=BLUE, text_color=WHITE, border_color=BLUE)
    add_box(sl, Cm(x_n), Cm(y_n + 0.55), Cm(16.0), Cm(2.4),
            body, font_size=10, bg_color=WHITE, border_color=BLUE2)


# ══════════════════════════════════════════════════════════════════
# PAGE 12 — 기계연 연구목표 및 수행 내용
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
add_header(sl, "한국기계연구원 연구목표 및 수행 내용 (2세부 공동기관)",
           "4개 항목 — 레시피개발 / 코팅층분석 / 레시피이전 / 전기화학평가")
add_page_num(sl, 12)

items_12 = [
    (
        "① ALD 공정 레시피 개발 (4종 산화물)",
        "목표: 배치식 ALD 장비 2대를 활용한 H₂O 기반 산화물 ALD 공정 레시피 체계적 확립\n"
        "• 배치식 ALD #1/#2 병렬 운용: TiO₂/Al₂O₃(1차), V₂O₅/ZnO(2차)\n"
        "• Self-limiting 반응 조건 정량 확인 — GPC vs. pulse 시간 포화 곡선\n"
        "• 독창 구조 개발: V₂O₅(내층)/TiO₂(외층) ALD 이중층 레시피 최초 확립",
        [["TiO₂","TTIP","H₂O","200–250°C","0.5–0.8 Å/cyc"],
         ["Al₂O₃","TMA","H₂O","150–200°C","0.8–1.2 Å/cyc"],
         ["V₂O₅","VTIP","H₂O","150–250°C","0.3–0.6 Å/cyc"],
         ["ZnO","DEZn","H₂O","100–200°C","1.5–2.0 Å/cyc"]],
        ["코팅 물질","전구체","반응제","탐색 온도","목표 GPC"]
    ),
    (
        "② 코팅층 원자 단위 종합 분석",
        "목표: HR-TEM/EELS/XPS 기반 ALD 코팅층 구조·조성·두께 정밀 분석\n"
        "• HR-TEM/EELS: 코팅층 두께 0.1 nm 해상도 + 원소 분포 nm 수준 매핑\n"
        "• XPS: 코팅층 조성 및 화학 결합 상태 (Ti 2p, Al 2p, V 2p, Zn 2p)\n"
        "• 배치식 vs. 연속식 코팅 품질 동등성 HRTEM 직접 비교 검증",
        [["HR-TEM (FIB 포함)","코팅층 두께·계면 원자 구조"],
         ["EELS","원소 분포 고분해능 매핑"],
         ["XPS","표면 조성·화학 결합 상태"],
         ["SEM-EDX","입자 형상·원소 면분석"],
         ["Raman / FTIR","구조 변화·잔류 전구체"]],
        ["분석 장비","분석 목적"]
    ),
]

y_top = [2.6, 10.8]
for idx, (title, body, tbl_data, tbl_hdr) in enumerate(items_12):
    y_t = y_top[idx]
    add_box(sl, Cm(0.4), Cm(y_t), Cm(33.0), Cm(0.6),
            title, font_size=11, bold=True,
            bg_color=BLUE, text_color=WHITE, border_color=BLUE)
    add_box(sl, Cm(0.4), Cm(y_t + 0.6), Cm(16.0), Cm(4.0),
            body, font_size=9, bg_color=WHITE, border_color=BLUE2)
    add_table(sl, Cm(16.8), Cm(y_t + 0.6), Cm(16.6), Cm(4.0),
              tbl_hdr, tbl_data, font_size=9)

# 항목 3, 4 간략
items_34 = [
    ("③ 레시피 이전 및 연속식 공정 피드백 지원",
     "• 레시피 패키지 작성 — 온도·pulse/purge 시간·GPC·공정 윈도우 문서화\n"
     "• 단계별 이전: FBR(1차) → OPERIO 10(2차) → 10 kg/h 양산(3차)\n"
     "• RTD 기반 배치→연속식 변환 방법론 최초 확립 (정량적 매핑 알고리즘)"),
    ("④ 코팅 양극재 전기화학 성능 검증",
     "• Half-cell 기반 4종 코팅재·두께별 비교 평가\n"
     "• 200사이클 후 Post-mortem 분석 — TEM/XPS 기반 코팅층 안정성 규명\n"
     "• 1Ah 파우치셀 해체 분석 (씨오알엔 협력)"),
]
y34 = 15.2
for i, (t, b) in enumerate(items_34):
    x34 = 0.4 + i * 16.8
    add_box(sl, Cm(x34), Cm(y34), Cm(16.0), Cm(0.6),
            t, font_size=10, bold=True,
            bg_color=BLUE, text_color=WHITE, border_color=BLUE)
    add_box(sl, Cm(x34), Cm(y34 + 0.6), Cm(16.0), Cm(2.8),
            b, font_size=9, bg_color=WHITE, border_color=BLUE2)


# ══════════════════════════════════════════════════════════════════
# PAGE 15 — 연구역량
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
add_header(sl, "한국기계연구원 연구역량 — ALD 분야 논문·특허 실적",
           "20년+ ALD 연구 이력 | Science 포함 SCI 100편+ | 피인용 2,000회+")
add_page_num(sl, 15)

add_box(sl, Cm(0.4), Cm(2.6), Cm(33.0), Cm(0.6),
        "■ 핵심 논문 실적 (ALD·배터리 계면 관련 5선)", font_size=11, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)
add_table(sl, Cm(0.4), Cm(3.2), Cm(33.0), Cm(5.5),
    ["#", "논문 제목", "저널", "피인용", "본 과제 연관성"],
    [
        ["1", "Greatly increased toughness of infiltrated spider silk",
         "Science", "503", "ALD 원천기술 — 기상 전구체 침투·코팅 이론적 기원"],
        ["2", "Accelerated Sulfur Evolution by TiS₂/TiO₂@MXene",
         "Adv. Funct. Mater.", "135", "TiO₂ ALD 전극 계면 제어 → 본 과제 TiO₂ 레시피 선행"],
        ["3", "Porous carbon with VC/V₂O₃₋ₓ for Li–S batteries",
         "Energy Storage Mater.", "65", "V 산화물 코팅 계면 화학 → V₂O₅ ALD 레시피 선행 지식"],
        ["4", "Taming Li Nucleation by ZnF₂ ALD Layer",
         "Advanced Science", "신규(2025)", "ZnF₂ ALD 배터리 계면 안정화 직접 실증"],
        ["5", "Ultrathin 90 nm ZnF₂ Layer for Zn Anode",
         "ACS Energy Lett.", "신규(2025)", "ALD 초박막 음극 보호 — 두께 제어 기술"],
    ],
    font_size=9, header_font_size=10
)

# ALD 연구 이력 타임라인
add_box(sl, Cm(0.4), Cm(9.0), Cm(33.0), Cm(0.6),
        "■ ALD 연구 이력 타임라인", font_size=11, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)
timeline = [
    ("2006–2009", "독일 Max Planck\n연구소", "ALD 원천기술 확립\nScience 2009 (피인용 503회)\nZnO, TiO₂ ALD on biological templates"),
    ("2010–2015", "한국기계연구원", "ALD 나노구조 및\n전극 소재 응용\ngraphene CVD 기반 나노 복합"),
    ("2020–2023", "배터리 전극\nALD 코팅 응용", "TiO₂/TiS₂ 계면\nV₂O₃₋ₓ 표면 코팅\nLi-S 전지 적용"),
    ("2024–2025", "ALD 배터리\n계면 특화", "ZnF₂ ALD Li/Zn 음극\n(Adv. Sci., ACS Energy Lett.)\nSIB 양극재 ALD 착수 → 본 과제"),
]
x_t = 0.4
for period, org, content in timeline:
    add_box(sl, Cm(x_t), Cm(9.65), Cm(8.0), Cm(0.7),
            period, font_size=9, bold=True,
            bg_color=NAVY, text_color=WHITE, border_color=NAVY, align=PP_ALIGN.CENTER)
    add_box(sl, Cm(x_t), Cm(10.35), Cm(8.0), Cm(0.55),
            org, font_size=8,
            bg_color=BLUE2, text_color=WHITE, border_color=BLUE2, align=PP_ALIGN.CENTER)
    add_box(sl, Cm(x_t), Cm(10.9), Cm(8.0), Cm(3.5),
            content, font_size=9, bg_color=WHITE, border_color=BLUE2)
    x_t += 8.4

# 역량 요약 박스 3종
add_box(sl, Cm(0.4), Cm(14.6), Cm(33.0), Cm(0.6),
        "■ 핵심 역량 요약", font_size=11, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)
caps = [
    ("ALD 연구 경력", "20년+ / Science 포함 SCI 논문 100편+\n누적 피인용 2,000회+"),
    ("장비 역량",     "배치식 ALD 장비 2대 보유 및 실운용\n4종 산화물 공정 가능"),
    ("분석 역량",     "HRTEM/EELS 기반 원자 단위\nALD 코팅층 구조 분석 기술 내재화"),
]
x_c = 0.4
for cap_title, cap_body in caps:
    add_box(sl, Cm(x_c), Cm(15.25), Cm(10.8), Cm(0.6),
            cap_title, font_size=10, bold=True,
            bg_color=NAVY, text_color=WHITE, border_color=NAVY, align=PP_ALIGN.CENTER)
    add_box(sl, Cm(x_c), Cm(15.85), Cm(10.8), Cm(2.8),
            cap_body, font_size=10, bg_color=RGBColor(0xE8,0xF4,0xFF), border_color=BLUE2)
    x_c += 11.2


# ══════════════════════════════════════════════════════════════════
# PAGE 24 — 기관별 핵심 보유 기술
# ══════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
set_bg(sl)
add_header(sl, "기관별 핵심 보유 기술 — 2세부 컨소시엄 핵심 강점",
           "한국기계연구원 (공정기술·분석역량) × 대성기계공업㈜ (장비기술·양산역량)")
add_page_num(sl, 24)

# 기계연 섹션
add_box(sl, Cm(0.4), Cm(2.6), Cm(16.0), Cm(0.65),
        "한국기계연구원 (KIMM)", font_size=12, bold=True,
        bg_color=BLUE, text_color=WHITE, border_color=BLUE)

kimm_items = [
    ("① ALD 공정 기술",
     "• 배치식 ALD 장비 2대: TMA/TTIP/VTIP/DEZn 전구체 취급\n"
     "• Self-limiting 반응 정밀 측정 기술 (GPC vs. pulse 포화 곡선)\n"
     "• 독창 구조: V₂O₅(내층)/TiO₂(외층) 이중층 ALD 레시피 개발 추진"),
    ("② 나노 분석 인프라",
     "• HR-TEM / EELS: 코팅층 두께 0.1 nm 해상도 + 원소 분포 매핑\n"
     "• XPS: 표면 조성·화학 결합 상태 정밀 분석\n"
     "• SEM-EDX / Raman / FTIR / Potentiostat 완비"),
    ("③ 배터리 계면 연구 실적",
     "• ZnF₂ ALD → 배터리 계면 안정화 (Adv. Sci. 2025)\n"
     "• TiO₂/TiS₂ 전극 계면 제어 (Adv. Funct. Mater. 2023, 피인용 135)\n"
     "• 바나듐 산화물 표면 코팅 (Energy Storage Mater. 2022, 피인용 65)"),
]
y_k = 3.25
for t, b in kimm_items:
    add_box(sl, Cm(0.4), Cm(y_k), Cm(16.0), Cm(0.55),
            t, font_size=10, bold=True,
            bg_color=BLUE2, text_color=WHITE, border_color=BLUE2)
    add_box(sl, Cm(0.4), Cm(y_k + 0.55), Cm(16.0), Cm(3.1),
            b, font_size=9, bg_color=WHITE, border_color=BLUE2)
    y_k += 3.7

# 대성기계 섹션
add_box(sl, Cm(17.0), Cm(2.6), Cm(16.4), Cm(0.65),
        "대성기계공업㈜", font_size=12, bold=True,
        bg_color=GREEN, text_color=WHITE, border_color=GREEN)

ds_items = [
    ("① 연속식 분말 ALD 장비 (OPERIO 10)",
     "• 공압 이송 기반 상압 연속 공정 — 국내 유일 독자 개발\n"
     "• 분말 체류시간(RTD) 정밀 제어 능력\n"
     "• 구간 분리형 반응기 구조 (전구체 교차 오염 방지)"),
    ("② 스케일업 실증 역량",
     "• R&D FBR: 소량 / 두께 편차 ≤10%  (1차년도)\n"
     "• OPERIO 10 (Pilot): kg급 / 두께 편차 ≤7%  (2차년도)\n"
     "• 양산형 연속식: 10 kg/h ★ / 두께 편차 ≤5% ★  (3차년도)"),
    ("③ Turn-key 공급 역량",
     "• ALD 장비 Turn-key 패키지 (레시피+장비+SOP) 상용화 경험\n"
     "• 국내 공압 이송 기반 연속식 분말 ALD 양산 장비 국내 선도\n"
     "• 배터리 소재 ALD 코팅 공정 국산화 → 해외 의존도 저감"),
]
y_d = 3.25
for t, b in ds_items:
    add_box(sl, Cm(17.0), Cm(y_d), Cm(16.4), Cm(0.55),
            t, font_size=10, bold=True,
            bg_color=RGBColor(0x1E,0x8B,0x50), text_color=WHITE,
            border_color=GREEN)
    add_box(sl, Cm(17.0), Cm(y_d + 0.55), Cm(16.4), Cm(3.1),
            b, font_size=9, bg_color=WHITE, border_color=GREEN)
    y_d += 3.7

# 하단 핵심 요약 배너
add_box(sl, Cm(0.4), Cm(17.1), Cm(33.0), Cm(1.7),
        "RTD 기반 배치→연속식 변환 방법론 최초 확립  |  국내 유일 배치식↔연속식 분말 ALD 협력 모델\n"
        "SIB 양극재 V₂O₅/TiO₂ ALD 이중층 적용 — 세계 최초 시도",
        font_size=12, bold=True,
        bg_color=NAVY, text_color=WHITE, border_color=NAVY, align=PP_ALIGN.CENTER)


# ── 저장
OUT = r"d:\23_GitHub\01_Office_PC\030_Project\032_Proposed Project\08_2026년 첨단소재 사업화 기술 개발(성과확장형)\발표자료_초안_8페이지.pptx"
prs.save(OUT)
print(f"[완료] 저장: {OUT}")
